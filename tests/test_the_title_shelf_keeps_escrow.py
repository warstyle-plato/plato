"""Плитка «На эскроу» доезжает до титульного слайда.

Владелец, 01.09.2026: «надо куда-то на 1 слайд остаток на эскроу». В отчёте эта
плитка есть и стоит шестой, а полка титула держала пять — шестая срезалась
молча. Отрезанная плитка неотличима от непосчитанной: на слайде её просто нет,
и почему — не сказано нигде.

Колода при этом по-прежнему ничего не считает: плитка берётся из той же
вёрстки, что печатает отчёт.

Запуск: python3 -m pytest tests/test_the_title_shelf_keeps_escrow.py -q
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from market_search import sales_deck  # noqa: E402

# Шапка свода: та же полка из шести плиток, что печатает отчёт.
SHELF = """
<section class="salesblock"><div class="blockhead"><h3>Свод</h3></div>
<div class="kv">
<div><div>Договоров</div><div>77</div><div>31,2 млн ₽ на договор</div></div>
<div><div>Квартиры</div><div>57 из 220</div><div>25,9% лотов проекта</div></div>
<div><div>Метры квартир</div><div>3 173 м²</div><div>23,6% из 13 429 м²</div></div>
<div><div>Машино-места</div><div>14 из 75</div><div>18,7% мест проекта</div></div>
<div><div>Выручка</div><div>2 400,9 млн ₽</div><div>18,3% из 13 134,7 млн ₽</div></div>
<div><div>На эскроу</div><div>943,5 млн ₽</div><div>39,3% от продаж</div></div>
</div>
<div class="sumup">Продано 18,3% ожидаемой выручки проекта.</div></section>
"""


def _first_slide_text() -> str:
    pytest.importorskip("pptx")
    from pptx import Presentation

    raw = sales_deck.build(sales_deck.sections(SHELF), title="Продажи — Проба",
                           subtitle="срез", footer="DevelopAid")
    slide = Presentation(io.BytesIO(raw)).slides[0]
    return " | ".join(
        run.text
        for shape in slide.shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs)


def test_escrow_is_on_the_title_slide() -> None:
    text = _first_slide_text()
    assert "На эскроу" in text, "шестая плитка срезалась — на титуле её нет"
    assert "943,5 млн ₽" in text
    # Остальные пять на месте: шестая добавлена, а не заменила соседку.
    for name in ("Договоров", "Квартиры", "Метры квартир", "Машино-места", "Выручка"):
        assert name in text, name


def test_the_deck_still_counts_nothing() -> None:
    """Плитка взята из вёрстки отчёта, а не посчитана заново."""
    body = (ROOT / "market_search" / "sales_deck.py").read_text(encoding="utf-8")
    start = body.index("    def put_shelf(")
    shelf = body[start:body.index("\n    def put_strip(", start)]
    for arithmetic in ("escrow", "/ 1e6", "* 100", "sum("):
        assert arithmetic not in shelf, (
            "полка только раскладывает числа отчёта, считать ей нечего")
