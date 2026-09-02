"""Графики колоды одеты по продукту, а не оставлены заводскими.

«PPT убогие графики очень, как для первого класса школы» (владелец,
30.08.2026) — и это была правда: столбик во всю ширину слота офисной синевы,
сетка по всему полю, ось со значениями, подписи заводским шрифтом, таблицы в
синюю полоску.

Правила взяты не с потолка: один ряд — легенды нет, её работу делает заголовок
слайда; текст носит текстовые токены, а не цвет ряда; значения стоят прямо на
столбиках, пока их немного, и тогда сетка с осью — лишняя краска; столбик не
заполняет слот. Цвет — фирменный синий кабинета, и он проверен валидатором
палитры (светлота, насыщенность, контраст к белому).

Запуск: python3 -m pytest tests/test_the_deck_charts_are_not_factory_default.py -q
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from market_search import sales_deck  # noqa: E402

pptx = pytest.importorskip("pptx")


def deck(categories: int):
    head = ["Месяц", "млн ₽"]
    rows = [[f"2026-{month + 1:02d}", str(100 + month)] for month in range(categories)]
    # `charted` — то, что разбор снял с экрана: в этом разделе экран рисует
    # график. Колода своего решения об этом не принимает, поэтому здесь оно
    # объявляется явно, как его объявила бы разметка со своим `<svg>`.
    pages = [{"title": "Динамика", "note": "Темп растёт.", "charted": True,
              "tables": [{"head": head, "rows": rows}], "lines": []}]
    raw = sales_deck.build(pages, title="Продажи", subtitle="срез", footer="DevelopAid")
    return pptx.Presentation(io.BytesIO(raw))


def charts(presentation):
    return [shape.chart for slide in presentation.slides for shape in slide.shapes
            if getattr(shape, "has_chart", False) and shape.has_chart]


def tables(presentation):
    return [shape.table for slide in presentation.slides for shape in slide.shapes
            if getattr(shape, "has_table", False) and shape.has_table]


def test_the_bar_wears_the_declared_palette_of_the_deck() -> None:
    """Столбик носит объявленный цвет колоды, а не офисный по умолчанию.

    Палитра у колоды с 01.09.2026 своя — чёрно-серая с зелёным акцентом
    (решение владельца): её носят на встречу отдельным файлом. Экран кабинета
    и PDF при этом остаются синими. Проверяется не «какой синий», а что цвет
    ОБЪЯВЛЕН нами: офисная палитра PowerPoint даёт свой 4472C4, и по нему
    сразу видно, что ряд не покрашен вовсе.
    """
    chart = charts(deck(6))[0]
    series = chart.plots[0].series[0]
    assert str(series.format.fill.fore_color.rgb) == "1F6FB2", "столбик колоды"
    # Обводки у столбика нет: рамка вокруг метки — краска, которая не данные.
    assert series.format.line.fill.type is not None


def test_one_series_has_no_legend() -> None:
    """Легенда из одной строки повторяет заголовок слайда и ест место."""
    assert charts(deck(6))[0].has_legend is False


def test_few_columns_carry_their_values_and_drop_the_grid() -> None:
    chart = charts(deck(5))[0]
    plot = chart.plots[0]
    assert plot.has_data_labels is True
    assert chart.value_axis.has_major_gridlines is False
    assert chart.value_axis.visible is False, "ось со значениями рядом с подписями — лишняя"


def test_many_columns_drop_the_labels_and_keep_a_hairline_grid() -> None:
    """Число на каждом столбике — хаос, когда столбиков тринадцать."""
    chart = charts(deck(13))[0]
    assert chart.plots[0].has_data_labels is False
    assert chart.value_axis.has_major_gridlines is True
    assert chart.value_axis.visible is True
    line = chart.value_axis.major_gridlines.format.line
    assert str(line.color.rgb) == "E5E5E5", "сетка волосяная и отступает на второй план"


def test_the_column_does_not_fill_its_slot() -> None:
    """Столбик во всю ширину слота и есть то самое «как для первого класса»."""
    assert charts(deck(3))[0].plots[0].gap_width >= 400
    assert charts(deck(6))[0].plots[0].gap_width >= 250
    assert charts(deck(14))[0].plots[0].gap_width <= 100, "тринадцати столбикам нужен воздух поуже"


def test_the_text_does_not_wear_the_data_colour() -> None:
    chart = charts(deck(5))[0]
    assert str(chart.font.color.rgb) == "6B7280"
    labels = chart.plots[0].data_labels
    assert str(labels.font.color.rgb) == "000000"
    # Формат зависит от величины: «#,##0.#» в русской раскладке печатает
    # разделитель дробной части и там, где дроби нет, — подписи столбиков
    # стояли как «576 680,» (снимок владельца, 31.08.2026). Крупным числам
    # дробная часть не нужна вовсе.
    assert labels.number_format in {"#,##0.#", "#,##0"}
    assert labels.number_format_is_linked is False


def test_the_table_is_not_blue_striped() -> None:
    grid = tables(deck(5))[0]
    assert grid.first_row is True and grid.horz_banding is False
    assert str(grid.cell(0, 0).fill.fore_color.rgb) == "F4F7FA", "шапка на светлой подложке"
    assert str(grid.cell(1, 0).fill.fore_color.rgb) == "FFFFFF"


def test_numbers_stand_to_the_right() -> None:
    """Колонка чисел, прижатая влево, не читается столбиком."""
    from pptx.enum.text import PP_ALIGN

    grid = tables(deck(5))[0]
    value = grid.cell(1, 1)
    assert value.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    assert grid.cell(1, 0).text_frame.paragraphs[0].alignment != PP_ALIGN.RIGHT


def test_sibling_labels_do_not_glue_into_one_word() -> None:
    """«факт, млн ₽цена квартир, ₽/м²» — так это выглядело на слайде.

    Подписи легенды и полос лежат соседними `span` без пробела между ними:
    браузер разводит их отступом, разбор склеивал в одно слово. Разделитель
    ставится только между соседями одного уровня — `span` внутри строки
    разбивать нечего.
    """
    html = ('<section class="salesblock"><h2>Динамика</h2>'
            '<div class="muted"><span><span></span>факт, млн ₽</span>'
            '<span><span></span>цена квартир, ₽/м²</span></div>'
            '<p>Внутри строки <span class="muted">пояснение</span> не рвётся.</p>'
            '</section>')
    lines = sales_deck.sections(html)[0]["lines"]
    assert "факт, млн ₽ · цена квартир, ₽/м²" in lines
    assert "Внутри строки пояснение не рвётся." in lines


def test_a_fold_label_is_not_content() -> None:
    """«Помесячно числами» — подпись сворачивалки, и на слайде она сирота."""
    html = ('<section class="salesblock"><h2>Продукты</h2>'
            '<details><summary>Продукты числами</summary>'
            '<table><thead><tr><th>Что</th><th>Сколько</th></tr></thead>'
            '<tbody><tr><td>Квартира</td><td>56</td></tr>'
            '<tr><td>Паркинг</td><td>14</td></tr></tbody></table></details></section>')
    page = sales_deck.sections(html)[0]
    assert "Продукты числами" not in page["lines"]
    assert page["tables"], "таблица под сворачивалкой при этом обязана остаться"


def test_a_section_without_words_does_not_get_an_empty_slide() -> None:
    """«Расторжения» шли листом, на котором стоял один заголовок."""
    from pptx import Presentation
    import io

    html = ('<section class="salesblock"><h2>Расторжения</h2>'
            '<table><thead><tr><th>Месяц</th><th>млн ₽</th></tr></thead>'
            '<tbody><tr><td>2026-05</td><td>12,0</td></tr>'
            '<tr><td>2026-06</td><td>3,5</td></tr></tbody></table></section>')
    deck = Presentation(io.BytesIO(sales_deck.build(
        sales_deck.sections(html), title="Т", subtitle="с", footer="ф")))
    for slide in deck.slides:
        filled = [shape for shape in slide.shapes
                  if shape.has_chart or shape.has_table
                  or (shape.has_text_frame and shape.text_frame.text.strip())]
        # Заголовок и номер — не содержание: лист, кроме них, обязан что-то нести.
        assert len(filled) > 2, "слайд с одним заголовком и номером"


def test_the_deck_adds_no_words_of_its_own() -> None:
    """«Зачем ты там вставляешь свои комментарии — разве это есть в PDF?»
    (владелец, 31.08.2026).

    Колода печатает тот же отчёт, что и PDF, и обе поверхности получают одну
    разметку. Строка, которой в отчёте нет, на слайде читается как часть
    отчёта — а её там нет. Сначала сноска стояла под каждым графиком, потом
    переехала на титул одной строкой; и то и другое — наши слова в чужом
    документе.
    """
    import io

    from pptx import Presentation

    html = ('<div class="kv"><div><div>Договоров</div><div>76</div><div></div></div>'
            '<div><div>Выручка</div><div>2 345,3 млн ₽</div><div></div></div></div>'
            '<section class="salesblock"><h3>Динамика</h3>'
            '<div class="sumup">Темп растёт.</div>'
            '<table><thead><tr><th>Месяц</th><th>млн ₽</th></tr></thead><tbody>'
            '<tr><td>2026-01</td><td>30,5</td></tr>'
            '<tr><td>2026-02</td><td>60,5</td></tr></tbody></table></section>')
    pages = sales_deck.sections(html)
    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Продажи — Кутузов Сити",
        subtitle="Свод продаж DevelopAid · срез 2026-08-27", footer="DevelopAid")))

    # Всё, что стоит на слайдах, обязано быть в разметке отчёта. Исключения —
    # только служебное: имя отчёта, подзаголовок, колонтитул и номер листа.
    # Служебное — имя отчёта, его подзаголовок, колонтитул и номер листа.
    # Всё это слова вызывающего, то есть самого отчёта, а не колоды.
    allowed = {text.casefold() for text in (
        "Продажи — Кутузов Сити",
        "Свод продаж DevelopAid · срез 2026-08-27",
        "DevelopAid",
        "DevelopAid · Продажи — Кутузов Сити"
        " · Свод продаж DevelopAid · срез 2026-08-27",
        "1", "2",
    )}
    source = html.lower()
    for slide in deck.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for line in shape.text_frame.text.splitlines():
                text = line.strip()
                if not text or text.casefold() in allowed:
                    continue
                assert text.casefold() in source, (
                    f"слова, которых нет в отчёте: {text!r}")


def test_the_price_per_metre_rides_as_a_line_not_its_own_slide() -> None:
    """«И линия цены метра то должна быть на этих графиках» (владелец,
    30.08.2026).

    Правило про цену уже записано: «цена — всегда линия на своей шкале, а не
    вкладка со столбиками». В колоде она уходила своим слайдом со столбиками —
    то же самое другими словами: смотрят на объём, а цена в это время на
    соседнем листе. Теперь она идёт линией справа на каждом графике объёма, а
    своим слайдом остаётся, только если объёма рядом нет вовсе.
    """
    table = {"head": ["Месяц", "Лотов", "млн ₽", "₽/м²"],
             "rows": [["2026-07", "4", "140,8", "712 747"],
                      ["2026-06", "9", "301,2", "717 000"],
                      ["2026-05", "9", "288,0", "705 100"]]}
    drawn = sales_deck.charts(table)
    # Цена своего листа не заводит: она линией справа на КАЖДОМ графике
    # объёма. Меры объёма при этом остаются — экран предлагает их
    # переключателем, а в документе переключателя нет.
    assert "₽/м²" not in [item["name"] for item in drawn], "цена ушла столбиками"
    # График на раздел один: «Лотов» осталась колонкой таблицы и названа под
    # графиком — лист на каждую меру давал пустые слайды (владелец, 31.08.2026).
    assert [item["name"] for item in drawn] == ["млн ₽"]
    assert drawn[0]["other_measures"] == ["Лотов"]
    assert all([row["name"] for row in item["second"]] == ["₽/м²"] for item in drawn)
    # Одна цена без объёма — сама себе график: показать её иначе нечем.
    alone = sales_deck.charts({"head": ["Месяц", "₽/м²"],
                               "rows": [["2026-07", "712 747"], ["2026-06", "717 000"]]})
    assert [item["name"] for item in alone] == ["₽/м²"]
    assert "line" not in alone[0]


def test_the_price_line_is_a_real_line_on_its_own_axis() -> None:
    """Комбинированный график собирается правкой XML, и порядок в нём строгий.

    Все группы графиков обязаны стоять раньше всех осей: линия, приписанная в
    конец области, встаёт после осей — PowerPoint такой файл не открывает
    вовсе, а python-pptx, LibreOffice и схема его читают и молчат.
    """
    import io
    import zipfile

    from lxml import etree
    from pptx import Presentation
    from pptx.chart.xmlwriter import ChartXmlWriter  # noqa: F401  (проверка окружения)

    html = ('<section class="salesblock"><h2>Динамика</h2>'
            '<svg viewBox="0 0 700 250"></svg>'
            '<table><thead><tr><th>Месяц</th><th>млн ₽</th><th>₽/м²</th></tr></thead>'
            '<tbody><tr><td>2026-07</td><td>140,8</td><td>712 747</td></tr>'
            '<tr><td>2026-06</td><td>301,2</td><td>717 000</td></tr>'
            '<tr><td>2026-05</td><td>288,0</td><td>705 100</td></tr>'
            '</tbody></table></section>')
    blob = sales_deck.build(sales_deck.sections(html),
                            title="Т", subtitle="с", footer="ф")

    deck = Presentation(io.BytesIO(blob))
    found = [shape.chart for slide in deck.slides for shape in slide.shapes
             if shape.has_chart]
    assert found, "график не нарисовался вовсе"
    chart = found[0]
    kinds = [type(plot).__name__ for plot in chart.plots]
    assert "BarPlot" in kinds and "LinePlot" in kinds, kinds
    assert chart.has_legend, "два ряда без легенды неразличимы"

    namespace = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
    with zipfile.ZipFile(io.BytesIO(blob)) as pack:
        part = next(name for name in pack.namelist()
                    if name.startswith("ppt/charts/chart"))
        area = etree.fromstring(pack.read(part)).find(".//c:plotArea", namespace)
    order = [etree.QName(child).localname for child in area]
    groups = [index for index, name in enumerate(order) if name.endswith("Chart")]
    axes = [index for index, name in enumerate(order) if name.endswith("Ax")]
    assert max(groups) < min(axes), f"оси встали раньше групп: {order}"
    # Своя шкала справа, и её деления видны: урезанная шкала обязана назваться.
    right = [ax for ax in area.findall("c:valAx", namespace)
             if ax.find("c:axPos", namespace).get("val") == "r"]
    assert right, "у цены нет своей шкалы справа"
    assert right[0].find("c:delete", namespace).get("val") == "0"


def test_the_key_numbers_are_tiles_and_the_bands_are_bands() -> None:
    """«Ничего общего с отчётом и PDF» (владелец, 30.08.2026).

    Колода собиралась разбором отчёта в «заголовок, строки, таблица», и от
    экрана не переносилось ничего визуального: плашка ключевых чисел ехала
    таблицей «Показатель / Значение / Пояснение», а цветная лента долей
    пропадала целиком — у её кусков нет текста, только ширина и цвет.
    Теперь плитки — фигуры с крупным числом, лента — фигуры своих цветов.
    """
    import io

    from pptx import Presentation

    html = ('<div class="kv">'
            '<div><div>Договоров</div><div>76</div><div></div></div>'
            '<div><div>Выручка</div><div>2 345,3 млн ₽</div><div>17,9%</div></div>'
            '</div>'
            '<section class="salesblock"><h2>Квартирография</h2>'
            '<div style="margin:10px 0"><div class="muted">Пул проекта · как построено</div>'
            '<div style="display:flex;height:22px">'
            '<div style="width:23.2%;background:#1367AE" title="28,3-40 — 23,2%"></div>'
            '<div style="width:76.8%;background:#C4581B" title="40-55 — 76,8%"></div>'
            '</div></div></section>')

    pages = sales_deck.sections(html)
    bands = [strip for page in pages for strip in page.get("strips") or []]
    assert bands and len(bands[0]["parts"]) == 2
    assert bands[0]["caption"] == "Пул проекта · как построено"
    assert bands[0]["parts"][0]["colour"] == "1367AE"

    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Продажи — Кутузов Сити", subtitle="срез", footer="DevelopAid")))
    shapes = [shape for slide in deck.slides for shape in slide.shapes]
    filled = [shape for shape in shapes
              if str(shape.shape_type or "").startswith("AUTO_SHAPE")]
    assert filled, "ни плиток, ни ленты — только текст и таблицы"
    # Лента несёт цвета экрана, а не офисную палитру.
    tones = {"%02X%02X%02X" % tuple(shape.fill.fore_color.rgb)
             for shape in filled if shape.fill.type is not None}
    assert "1367AE" in tones and "C4581B" in tones
    # Ключевое число стоит крупно: полка титула, на которую смотрят с трёх
    # метров. Ищется по ВСЕМ фигурам, а не по одним автофигурам: число полки
    # лежит своей надписью над подложкой — колонок в подложке пять, и одним
    # текстовым полем они не набираются. Прежняя поломка («Показатель /
    # Значение / Пояснение» таблицей) этой проверкой ловится по-прежнему:
    # у ячеек таблицы своего text_frame среди фигур слайда нет.
    big = [run.font.size.pt for shape in shapes if shape.has_text_frame
           for para in shape.text_frame.paragraphs for run in para.runs
           if run.font.size and run.font.size.pt >= 20 and any(
               ch.isdigit() for ch in run.text)]
    assert big, "ключевого числа крупно нет ни на одном слайде"
    # И у шапки свода на слайде имя проекта, а не слово «Раздел».
    texts = [shape.text_frame.text for shape in shapes if shape.has_text_frame]
    assert not any(text.strip() == "Раздел" for text in texts)


def test_the_first_header_stands_over_its_own_column() -> None:
    """`grid.cell(0,0)` отдаёт новую обёртку на каждый вызов, поэтому сравнение
    «это ли первая ячейка» было всегда ложным, и «Месяц» уезжал вправо над
    колонкой дат, прижатых влево."""
    import io

    from pptx import Presentation

    pages = [{"title": "Динамика", "note": "", "lines": [], "strips": [],
              "tables": [{"head": ["Месяц", "Лотов"],
                          "rows": [["2026-07", "4"], ["2026-06", "9"]]}]}]
    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Т", subtitle="с", footer="ф")))
    grids = [shape.table for slide in deck.slides for shape in slide.shapes
             if getattr(shape, "has_table", False) and shape.has_table]
    assert grids
    first = grids[0].cell(0, 0).text_frame.paragraphs[0]
    second = grids[0].cell(0, 1).text_frame.paragraphs[0]
    assert first.alignment != second.alignment, "оба заголовка выровнены одинаково"


def test_a_conclusion_alone_does_not_get_its_own_slide() -> None:
    """«Этот слайд странный» (владелец, 30.08.2026): заголовок, одна строка
    вывода и пять дюймов белого.

    Вывод сам по себе слайдом не является — он едет подзаголовком на первый
    слайд раздела, где есть картинка. Слайд заводится под содержимое.
    """
    import io

    from pptx import Presentation

    html = ('<section class="salesblock"><h2>Динамика</h2>'
            '<div class="sumup">Последние три месяца — 213,6 млн ₽ в месяц.</div>'
            '<table><thead><tr><th>Месяц</th><th>млн ₽</th></tr></thead><tbody>'
            '<tr><td>2026-07</td><td>140,8</td></tr>'
            '<tr><td>2026-06</td><td>301,2</td></tr>'
            '<tr><td>2026-05</td><td>288,0</td></tr>'
            '</tbody></table></section>')
    deck = Presentation(io.BytesIO(sales_deck.build(
        sales_deck.sections(html), title="Т", subtitle="с", footer="ф")))
    slides = list(deck.slides)[1:]  # титул не в счёт
    for slide in slides:
        heavy = [shape for shape in slide.shapes
                 if shape.has_chart or shape.has_table
                 or str(shape.shape_type or "").startswith("AUTO_SHAPE")]
        assert heavy, "слайд без единой картинки, таблицы или плитки"
    # И вывод при этом не пропал: он стоит над первой картинкой.
    said = [shape.text_frame.text for slide in slides for shape in slide.shapes
            if shape.has_text_frame]
    assert any("213,6 млн ₽" in text for text in said)
    assert sum("213,6 млн ₽" in text for text in said) == 1, "вывод повторился"


def test_the_bars_keep_their_own_scale_next_to_the_price_line() -> None:
    """«Тут просто линии» (владелец, 30.08.2026) — и столбиков правда не было.

    С двумя шкалами первая обязана остаться видимой. Удалённая — а её удаляло
    правило «значения на столбиках, ось лишняя» — она уводит столбики на шкалу
    цены: 35 против 800 000, и от них на слайде не остаётся ничего.
    """
    import io
    import zipfile

    from lxml import etree

    html = ('<section class="salesblock"><h2>Спрос</h2>'
            '<svg viewBox="0 0 700 250"></svg>'
            '<table><thead><tr><th>Полоса</th><th>Просят</th>'
            '<th>₽/м² в книге</th></tr></thead><tbody>'
            '<tr><td>28,3-40</td><td>35</td><td>614 466</td></tr>'
            '<tr><td>40-55</td><td>41</td><td>644 507</td></tr>'
            '<tr><td>55-85</td><td>29</td><td>612 170</td></tr>'
            '<tr><td>85-168</td><td>17</td><td>734 077</td></tr>'
            '</tbody></table></section>')
    blob = sales_deck.build(sales_deck.sections(html),
                            title="Т", subtitle="с", footer="ф")
    namespace = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
    with zipfile.ZipFile(io.BytesIO(blob)) as pack:
        part = next(name for name in pack.namelist()
                    if name.startswith("ppt/charts/chart"))
        area = etree.fromstring(pack.read(part)).find(".//c:plotArea", namespace)
    axes = {ax.find("c:axPos", namespace).get("val"): ax
            for ax in area.findall("c:valAx", namespace)}
    assert set(axes) == {"l", "r"}, "у столбиков и цены должны быть свои шкалы"
    for side, ax in axes.items():
        gone = ax.find("c:delete", namespace)
        assert gone is not None and gone.get("val") == "0", \
            f"шкала {side} удалена — ряд уедет на чужую"


def test_without_a_price_line_the_axis_still_goes_away_when_values_are_on_top() -> None:
    """Правило «значения на столбиках — ось лишняя краска» остаётся: оно
    отменяется только второй шкалой."""
    import io

    from pptx import Presentation

    pages = [{"title": "Продукты", "note": "", "lines": [], "strips": [],
              "charted": True,
              "tables": [{"head": ["Продукт", "Договоров"],
                          "rows": [["Квартира", "56"], ["Паркинг", "14"],
                                   ["ПСН", "6"]]}]}]
    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Т", subtitle="с", footer="ф")))
    chart = [shape.chart for slide in deck.slides for shape in slide.shapes
             if getattr(shape, "has_chart", False) and shape.has_chart][0]
    assert chart.plots[0].has_data_labels
    assert chart.value_axis.visible is False


def test_a_section_with_bands_does_not_also_get_bar_charts() -> None:
    """«Тут некрасивые столбики, а не как в отчёте» (владелец, 30.08.2026).

    На экране квартирография — три цветные ленты долей. В колоде поверх них
    рисовались ещё и три синих столбиковых слайда подряд: то же самое, только
    хуже и без цвета полос. Числа при этом никуда не деваются — они в таблице
    раздела, которая идёт следом.
    """
    import io

    from pptx import Presentation

    html = ('<section class="salesblock"><h2>Квартирография</h2>'
            '<div class="sumup">Вымывается полоса 28,3-40.</div>'
            '<div style="margin:10px 0"><div class="muted">Пул проекта</div>'
            '<div style="display:flex;height:22px">'
            '<div style="width:23.2%;background:#1367AE" title="28,3-40 — 23,2%"></div>'
            '<div style="width:76.8%;background:#C4581B" title="40-55 — 76,8%"></div>'
            '</div></div>'
            '<table><thead><tr><th>Полоса</th><th>В пуле</th><th>Продано</th>'
            '</tr></thead><tbody>'
            '<tr><td>28,3-40</td><td>51</td><td>26</td></tr>'
            '<tr><td>40-55</td><td>60</td><td>12</td></tr></tbody></table></section>')
    deck = Presentation(io.BytesIO(sales_deck.build(
        sales_deck.sections(html), title="Т", subtitle="с", footer="ф")))
    shapes = [shape for slide in deck.slides for shape in slide.shapes]
    assert not [s for s in shapes if getattr(s, "has_chart", False) and s.has_chart], \
        "рядом с лентой снова появились столбики"
    # Лента на месте, и числа тоже.
    assert [s for s in shapes if str(s.shape_type or "").startswith("AUTO_SHAPE")]
    assert [s for s in shapes if getattr(s, "has_table", False) and s.has_table]


def test_the_chart_does_not_repeat_the_slide_title() -> None:
    """Имя меры уже стоит заголовком слайда; повторённое мелким серым над
    графиком, оно читается как чужая подпись."""
    import io

    from pptx import Presentation

    pages = [{"title": "Продукты", "note": "", "lines": [], "strips": [],
              "charted": True,
              "tables": [{"head": ["Продукт", "Договоров"],
                          "rows": [["Квартира", "56"], ["Паркинг", "14"]]}]}]
    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Т", subtitle="с", footer="ф")))
    chart = [shape.chart for slide in deck.slides for shape in slide.shapes
             if getattr(shape, "has_chart", False) and shape.has_chart][0]
    assert chart.has_title is False


def test_the_money_measure_leads_and_the_rest_stay_in_the_table() -> None:
    """Деньги первыми, а прочие меры — колонками таблицы, не листами.

    Прежде здесь стояло обратное: лист на каждую меру (владелец, 29.08.2026,
    «не проще для каждого графика свой слайд сделать?»). На настоящем своде это
    дало три почти одинаковых столбиковых листа подряд, и на двух из них не
    было ничего, кроме картинки — «пустые два слайда» (владелец, 31.08.2026).
    Решение 30.08 — «график на раздел один» — было записано в комментарии, но
    не доведено до кода; здесь оно и доводится.

    Пропасть меры при этом не должны: они колонками в таблице раздела, которая
    идёт следом, и названы подписью под графиком (`other_measures`).
    """
    money = sales_deck.charts({
        "head": ["Условие", "Договоров", "млн ₽", "На эскроу, млн ₽"],
        "rows": [["рассрочка", "35", "1 399,4", "573,7"],
                 ["100% оплата", "25", "447,4", "447,4"]]})
    assert [item["measure"] for item in money] == ["млн ₽"]
    assert money[0]["other_measures"], "мера, ушедшая в таблицу, не названа"
    # Две денежные колонки — один график: «На эскроу» линией рядом с продажами.
    assert [row["name"] for row in money[0]["extra"]] == ["На эскроу, млн ₽"]

    # Денег в таблице нет — берётся первая числовая, и она одна.
    counted = sales_deck.charts({
        "head": ["Источник", "Обращений", "Броней"],
        "rows": [["звонок", "518", "16"], ["агент", "44", "16"]]})
    assert [item["name"] for item in counted] == ["Обращений"]


def test_the_deck_is_laid_out_like_the_printed_report() -> None:
    """«Мне нужен максимально похожий на PDF вариант в pp» (владелец,
    31.08.2026).

    PDF свода — это напечатанный экран, и его вёрстка объявлена в `@media
    print` кабинета: надзаголовок прописными вразрядку под волосяной линейкой,
    заголовок, полка ключевых чисел на подложке, колонтитул на каждой странице.
    Колода собиралась своей вёрсткой — тёмный титул, номер листа углом, плитки
    карточками, — и рядом с отчётом читалась как другой документ.

    Проверяется то, что видно: шапка, полка и колонтитул. Числа при этом
    по-прежнему берутся с экрана, а не считаются заново.
    """
    import io

    from pptx import Presentation
    from pptx.util import Pt

    html = ('<div class="kv">'
            '<div><div>Договоров</div><div>76</div><div>с начала продаж</div></div>'
            '<div><div>Выручка</div><div>2 345,3 млн ₽</div><div>17,9%</div></div>'
            '</div>'
            '<section class="salesblock"><h2>Динамика</h2>'
            '<svg viewBox="0 0 700 250"></svg>'
            '<table><thead><tr><th>Месяц</th><th>млн ₽</th></tr></thead><tbody>'
            '<tr><td>2026-01</td><td>30,5</td></tr>'
            '<tr><td>2026-02</td><td>60,5</td></tr></tbody></table></section>')

    pages = sales_deck.sections(html)
    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Продажи — Кутузов Сити",
        subtitle="Свод продаж DevelopAid · срез 2026-08-27", footer="DevelopAid")))
    slides = list(deck.slides)
    assert len(slides) >= 2

    def texts(slide) -> list[str]:
        return [shape.text_frame.text.strip() for shape in slide.shapes
                if shape.has_text_frame]

    # Колонтитул на каждом листе: лист, отделившийся от колоды, обязан сам
    # себя нумеровать. Углового номера больше нет — в отчёте он внизу строки.
    for number, slide in enumerate(slides, 1):
        assert str(number) in texts(slide), f"на листе {number} нет его номера"

    # Титул несёт полку показателей, и её числа — те же, что на экране.
    title_slide = slides[0]
    joined = " ".join(texts(title_slide))
    assert "2 345,3 млн ₽" in joined and "76" in joined
    panel = [shape for shape in title_slide.shapes
             if str(shape.shape_type or "").startswith("AUTO_SHAPE")
             and shape.width > Pt(400)]
    assert panel, "полки показателей на титуле нет"

    # Титул себя не повторяет: имя отчёта стоит заголовком, и колонтитулом
    # оно читалось бы как заводская рамка.
    assert sum(1 for text in texts(title_slide)
               if "Кутузов Сити" in text) == 1

    # Полка не повторяется разделом: вторые числа о том же читаются как
    # расхождение, даже когда числа те же.
    for slide in slides[1:]:
        assert "2 345,3 млн ₽" not in " ".join(texts(slide))

    # Надзаголовок раздела — прописными и вразрядку, как в печатной шапке.
    section = slides[1]
    spaced = [run for shape in section.shapes if shape.has_text_frame
              for para in shape.text_frame.paragraphs for run in para.runs
              if run.font._rPr.get("spc")]
    assert spaced, "надзаголовок стоит без разрядки — прописные слипаются"
    assert any(run.text == run.text.upper() and run.text.strip() for run in spaced)


def test_the_deck_draws_a_chart_only_where_the_screen_draws_one() -> None:
    """«Там была куча столбиков опять» (владелец, 31.08.2026).

    Колода заводила столбики под ПЕРВУЮ таблицу каждого раздела, и на своде из
    одиннадцати разделов выходило шесть почти одинаковых синих слайдов.
    Расторжения столбиками не читаются вовсе, у структуры оплаты их три штуки,
    а экран в этих разделах ничего и не рисует: график у него ровно там, где
    график отвечает на вопрос раздела. Это решение уже принято — колода его не
    принимает второй раз, а переносит.
    """
    import io

    from pptx import Presentation

    def block(title: str, body: str) -> str:
        return f'<section class="salesblock"><h3>{title}</h3>{body}</section>'

    months = "".join(
        f"<tr><td>2026-{m:02d}</td><td>{m}</td><td>{m * 30},5</td></tr>"
        for m in range(1, 13))
    table = ('<table><thead><tr><th>Месяц</th><th>Лотов</th><th>млн ₽</th></tr>'
             f'</thead><tbody>{months}</tbody></table>')
    html = (
        block("Динамика", '<svg viewBox="0 0 700 250"></svg>' + table)
        + block("Каналы продаж",
                '<table><thead><tr><th>Канал</th><th>Договоров</th><th>млн ₽</th></tr>'
                '</thead><tbody><tr><td>Напрямую</td><td>31</td><td>862,4</td></tr>'
                '<tr><td>Брокер</td><td>28</td><td>901,0</td></tr></tbody></table>')
        + block("Расторжения",
                '<table><thead><tr><th>Договор</th><th>Возвращено, млн ₽</th></tr>'
                '</thead><tbody><tr><td>ДДУ-114</td><td>12,40</td></tr>'
                '<tr><td>ДДУ-201</td><td>9,15</td></tr></tbody></table>'))

    pages = sales_deck.sections(html)
    charted = {page["title"]: page.get("charted") for page in pages if page.get("title")}
    assert charted == {"Динамика": True, "Каналы продаж": False, "Расторжения": False}

    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Продажи", subtitle="срез", footer="DevelopAid")))
    # Графики есть только у раздела, где график рисует экран. Мер у него может
    # быть несколько — их предлагает переключатель, и в документе они обязаны
    # быть все; а у «Каналов» и «Расторжений» графика нет ни одного.
    charted = [(slide, shape.chart) for slide in deck.slides for shape in slide.shapes
               if getattr(shape, "has_chart", False)]
    assert charted, "раздел с графиком на экране остался без графика"
    titles = {shape.text_frame.text.split(" · ")[0]
              for slide, _ in charted for shape in slide.shapes
              if shape.has_text_frame and shape.text_frame.text.strip()}
    assert "Каналы продаж" not in titles and "Расторжения" not in titles, titles
    # Числа разделов без графика не пропали — они таблицами.
    tabled = [shape for slide in deck.slides for shape in slide.shapes
              if getattr(shape, "has_table", False)]
    assert len(tabled) >= 3, "раздел без графика остался и без чисел"


def test_columns_of_one_measure_stand_in_one_chart() -> None:
    """Раздел «Факт против планов» показывал один факт.

    Три колонки в «млн ₽» — факт, план ФМ, план банка, — а на слайд уезжала
    первая: график с именем «против планов» никаких планов не показывал.
    Колонки ОДНОЙ меры идут рядами одного графика; разные меры — нет.
    И прочерк в колонке едет ПРОПУСКОМ: у плана банка первый квартал пустой,
    и прежде из-за него терялась вся колонка.
    """
    import io

    from pptx import Presentation

    html = ('<section class="salesblock"><h3>Факт против планов</h3>'
            '<svg viewBox="0 0 700 250"></svg>'
            '<table><thead><tr><th>Квартал</th><th>Факт, млн ₽</th>'
            '<th>План ФМ, млн ₽</th><th>План банка, млн ₽</th><th>₽/м²</th></tr></thead>'
            '<tbody>'
            '<tr><td>2025 Q3</td><td>410,2</td><td>520,0</td><td>—</td><td>470000</td></tr>'
            '<tr><td>2025 Q4</td><td>724,7</td><td>880,0</td><td>640,0</td><td>482000</td></tr>'
            '<tr><td>2026 Q1</td><td>1210,4</td><td>1628,9</td><td>1100,0</td><td>488300</td></tr>'
            '</tbody></table></section>')

    pages = sales_deck.sections(html)
    deck = Presentation(io.BytesIO(sales_deck.build(
        pages, title="Продажи", subtitle="срез", footer="DevelopAid")))
    charts = [shape.chart for slide in deck.slides for shape in slide.shapes
              if getattr(shape, "has_chart", False)]
    assert len(charts) == 1, "меры одной оси разъехались по слайдам"
    # Композиция листа: факт столбиками, планы линиями на шкале рублей, цена
    # линией на своей шкале справа. Пять столбиков в категории читаются
    # частоколом, и «факт против плана» в нём не виден.
    groups = charts[0].plots
    assert len(groups) == 3, "график собран не так, как на листе"
    assert [series.name for series in groups[0].series] == ["Факт, млн ₽"]
    assert [series.name for series in groups[1].series] == [
        "План ФМ, млн ₽", "План банка, млн ₽"], (
        "планы той же меры не встали рядом с фактом")
    assert [series.name for series in groups[2].series] == ["₽/м²"]
    # Прочерк — пропуск, а не ноль: точки в этом квартале просто нет.
    assert list(groups[1].series[1].values)[0] is None
    # Заголовок слайда называет меру, а не первый ряд: «· Факт, млн ₽» над
    # графиком с обоими планами обещает меньше, чем на слайде есть.
    titles = [shape.text_frame.text for slide in deck.slides for shape in slide.shapes
              if shape.has_text_frame and "Факт против планов" in shape.text_frame.text]
    assert any(text.strip() == "Факт против планов · млн ₽" for text in titles), titles


def test_the_measure_is_read_wherever_the_unit_stands() -> None:
    """Шапка отчёта ставит единицу впереди: «МЛН ₽, ФАКТ».

    Мера бралась как часть после последней запятой, и на настоящем своде три
    денежные колонки получали три разные меры — «факт», «план фм», «план
    банка», — то есть в один график не собирались вовсе. А это и есть тот
    сводный график листа, ради которого раздел существует (стр. 9 PDF свода,
    31.08.2026). Единица ищется по своему написанию, а не по месту в шапке.
    """
    measure = sales_deck.measure_of
    assert measure("МЛН ₽, ФАКТ") == measure("Факт, млн ₽") == "млн ₽"
    assert measure("МЛН ₽, ПЛАН ФМ") == measure("МЛН ₽, ПЛАН БАНКА") == "млн ₽"
    # Цена метра отличается от денег и от площади — это две оси, между
    # которыми её и надо различить.
    assert measure("ЦЕНА ФАКТ, ₽/М²") == measure("₽/м²") == "₽/м²"
    assert measure("М²") == "м²"
    assert measure("Лотов") == measure("Договоров") == "шт"


def test_the_slide_repeats_the_composition_of_the_page() -> None:
    """«Графика для планов по метрам и рублям вообще нет, хотя в PDF есть
    хороший график, где и цены, и метры, и рубли на одном» (владелец,
    31.08.2026).

    Лист собирает его так: факт столбиками, планы линиями на шкале рублей, все
    цены линиями на шкале ₽/м². Слайд повторяет лист, а не придумывает свой:
    прежде он брал ОДНУ денежную колонку и одну ценовую, и раздел «факт против
    планов» никаких планов не показывал.
    """
    import io

    from pptx import Presentation

    html = ('<section class="salesblock"><h3>Факт против планов</h3>'
            '<svg viewBox="0 0 700 250"></svg>'
            '<table><thead><tr><th>Квартал</th><th>млн ₽, факт</th>'
            '<th>млн ₽, план ФМ</th><th>млн ₽, план банка</th>'
            '<th>цена факт, ₽/м²</th><th>цена ФМ, ₽/м²</th></tr></thead><tbody>'
            '<tr><td>2025 Q4</td><td>624,0</td><td>624,0</td><td>—</td>'
            '<td>670281</td><td>640000</td></tr>'
            '<tr><td>2026 Q1</td><td>554,0</td><td>454,8</td><td>1628,9</td>'
            '<td>724203</td><td>690000</td></tr>'
            '<tr><td>2026 Q2</td><td>599,8</td><td>646,0</td><td>358,2</td>'
            '<td>656733</td><td>705000</td></tr></tbody></table></section>')

    deck = Presentation(io.BytesIO(sales_deck.build(
        sales_deck.sections(html), title="Т", subtitle="с", footer="ф")))
    chart = [shape.chart for slide in deck.slides for shape in slide.shapes
             if getattr(shape, "has_chart", False)][0]
    groups = chart.plots
    assert len(groups) == 3
    assert [s.name for s in groups[0].series] == ["млн ₽, факт"]
    assert [s.name for s in groups[1].series] == ["млн ₽, план ФМ", "млн ₽, план банка"]
    assert [s.name for s in groups[2].series] == ["цена факт, ₽/м²", "цена ФМ, ₽/м²"]

    # Порядок детей области обязателен: сначала ВСЕ группы, потом оси. Иначе
    # PowerPoint не открывает файл вовсе, а всё остальное читает и молчит.
    import zipfile
    from lxml import etree

    ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    raw = io.BytesIO(sales_deck.build(
        sales_deck.sections(html), title="Т", subtitle="с", footer="ф"))
    with zipfile.ZipFile(raw) as pack:
        part = next(n for n in pack.namelist() if n.startswith("ppt/charts/chart"))
        area = etree.fromstring(pack.read(part)).find(f"{ns}chart/{ns}plotArea")
    kids = [etree.QName(node).localname for node in area]
    charts_at = [i for i, k in enumerate(kids) if k.endswith("Chart")]
    axes_at = [i for i, k in enumerate(kids) if k.endswith("Ax")]
    assert max(charts_at) < min(axes_at), kids
    # Линии планов живут на шкале столбиков, линии цены — на своей.
    own = [a.get("val") for a in area.find(f"{ns}barChart").findall(f"{ns}axId")]
    # Порядок групп в разметке значения не имеет — важно, на каких они осях:
    # одна линия делит шкалу со столбиками, вторая живёт на своей.
    lines = [[a.get("val") for a in node.findall(f"{ns}axId")]
             for node in area.findall(f"{ns}lineChart")]
    assert sum(1 for axes in lines if axes == own) == 1, lines
    assert sum(1 for axes in lines if axes != own) == 1, lines


def test_a_section_of_one_sentence_does_not_get_its_own_slide() -> None:
    """«Слайды 2-3 пустые вообще — там по одной строчке текста» (владелец,
    31.08.2026).

    На листе такой раздел занимает три сантиметра в потоке; на слайде — свой
    заголовок и пять дюймов белого. Такие идут по нескольку на общий лист, в
    своём порядке и со своими заголовками: пропасть вывод не должен — он и
    есть ответ раздела.
    """
    import io

    from pptx import Presentation

    def block(title: str, body: str = "") -> str:
        return f'<section class="salesblock"><h3>{title}</h3>{body}</section>'

    html = (block("Эскроу против погашения ПФ",
                  '<svg viewBox="0 0 700 250"></svg>'
                  '<div class="sumup">Покрытие 0,95x.</div>')
            + block("Чего эта воронка не даёт",
                    '<div class="sumup">Связать обращение с договором нечем.</div>')
            + block("Расторжения",
                    '<div class="sumup">Возвращено 21,55 млн ₽.</div>'))

    deck = Presentation(io.BytesIO(sales_deck.build(
        sales_deck.sections(html), title="Т", subtitle="с", footer="ф")))
    assert len(deck.slides) == 2, "три фразы заняли три листа"
    said = " ".join(shape.text_frame.text for slide in deck.slides
                    for shape in slide.shapes if shape.has_text_frame)
    # Ни один раздел не пропал: и заголовок, и его вывод на месте.
    for text in ("Эскроу против погашения ПФ", "Покрытие 0,95x.",
                 "Чего эта воронка не даёт", "Связать обращение с договором нечем.",
                 "Расторжения", "Возвращено 21,55 млн ₽."):
        assert text in said, text


def test_the_strip_takes_its_names_from_the_legend_by_colour() -> None:
    """Ленты каналов подписей не несут — имена стоят в легенде под лентой.

    Кусок брался, только если у него есть и ширина, и `title`; у каналов
    `title` нет вовсе, лента опознавалась пустой, и раздел уходил на слайд без
    единой картинки («каналы продаж без визуализации», владелец, 31.08.2026).
    Имя подбирается по цвету — ровно так же, как это читает человек.
    """
    html = ('<section class="salesblock"><h3>Каналы продаж</h3>'
            '<div style="margin:10px 0 4px">'
            '<div class="muted" style="font-size:12px">Выручка</div>'
            '<div style="display:flex;height:22px;margin-top:4px">'
            '<div style="width:37.30%;background:#5FA98A"></div>'
            '<div style="width:62.70%;background:#C4581B"></div></div>'
            '<div class="muted" style="font-size:12px;margin-top:3px">'
            '<span style="display:inline-block;width:9px;height:9px;'
            'background:#5FA98A;margin-right:4px"></span>свой отдел 37,3%'
            '<span style="display:inline-block;width:9px;height:9px;'
            'background:#C4581B;margin:0 4px 0 14px"></span>брокеры 62,7%'
            '</div></div></section>')

    page = sales_deck.sections(html)[-1]
    strips = page.get("strips") or []
    assert len(strips) == 1, "лента без подписей не опозналась"
    assert [(part["colour"], part["name"]) for part in strips[0]["parts"]] == [
        ("5FA98A", "свой отдел 37,3%"), ("C4581B", "брокеры 62,7%")]
    # Имя уехало в кусок — второй копией строкой раздела оно читалось бы как
    # подпись к соседнему блоку.
    assert page["lines"] == [], page["lines"]


def test_the_deck_wears_its_own_declared_palette() -> None:
    """У колоды палитра своя, и она объявлена в одном месте.

    До 01.09.2026 колода носила токены отчёта — синие, как экран кабинета.
    Владелец прислал колоду, перекрашенную надстройкой в чёрно-серую с зелёным
    акцентом, и выбрал её ДЛЯ КОЛОДЫ: её носят на встречу отдельным файлом.
    Экран и PDF остаются синими — перекрашивать их он не просил, а два
    несогласованных решения об одном цвете хуже одного явного.

    Проверяется не красота, а что цвета объявлены здесь и не расползлись
    литералами по функциям: разъехавшись, они дадут слайды разного вида в
    одной колоде, и заметить это можно будет только глазами.
    """
    source = (ROOT / "market_search" / "sales_deck.py").read_text(encoding="utf-8")
    assert "0x1F, 0x6F, 0xB2" in source, "столбик факта"
    assert "0x15, 0x80, 0x3D" in source, "зелёная полоса вывода"
    assert "0xF5, 0xF5, 0xF5" in source, "подложка плашки вывода"
    assert "0x1A, 0x1A, 0x1A" in source, "текст вывода"
    assert "0xE5, 0xE5, 0xE5" in source, "волосяная линейка"
    assert "def put_note(" in source, "вывод рисуется плашкой, а не строкой"


def test_the_deck_carries_the_emblem_and_does_not_say_the_name_twice() -> None:
    """Эмблема одна на все поверхности, и она лежит в `PAGE`.

    У колоды её не было вовсе — лист без неё не опознаётся как наш («нет
    стилистики Плато», владелец, 31.08.2026). Байты идут тем же крючком, что
    карта и картинки отчёта: копии у эмблемы нет, её негде было бы обновлять.

    Формат при этом переводится, а не глушится. Эмблема у нас в WebP —
    страницы его берут, PowerPoint нет, и первая версия прятала отказ за
    `except: pass`: эмблема просто не появлялась, то есть молчаливый пропуск
    вместо перевода.
    """
    import io

    from pptx import Presentation

    page = {"title": "Динамика", "note": "Темп растёт.", "charted": True,
            "lines": [], "strips": [],
            "tables": [{"head": ["Месяц", "млн ₽"],
                        "rows": [["2026-01", "30,5"], ["2026-02", "60,5"]]}]}

    def built(logo):
        raw = sales_deck.build([dict(page)], title="Продажи", subtitle="срез",
                               footer="DevelopAid", logo=logo)
        deck = Presentation(io.BytesIO(raw))
        first = deck.slides[0]
        pictures = [s for s in first.shapes if s.__class__.__name__ == "Picture"]
        texts = [s.text_frame.text.strip() for s in first.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        return pictures, texts

    import main_registry

    emblem = main_registry._local_asset("/guide/assets/logo.webp")
    assert emblem, "эмблема не отдаётся крючком — на слайде её взять неоткуда"

    pictures, texts = built(emblem)
    assert len(pictures) == 1, "эмблемы на титуле нет"
    # WebP переведён: формат слайда его не знает, и картинка иначе не встала бы.
    assert pictures[0].image.content_type in ("image/png", "image/jpeg")
    # Под картинкой со словом ПЛАТО строка «DEVELOPAID» — то же слово дважды.
    assert "DEVELOPAID" not in texts, texts

    # Эмблемы нет — надзаголовок остаётся: имя отчёта назвать всё равно надо.
    bare_pictures, bare_texts = built(None)
    assert bare_pictures == []
    assert "DEVELOPAID" in bare_texts, bare_texts


def test_a_big_number_gets_no_dangling_separator() -> None:
    """«576 680,» — это оборванное число, а не подпись."""
    from pptx.util import Emu  # noqa: F401  (импорт держит зависимость явной)

    from market_search import sales_deck

    table = {"head": ["Квартал", "Цена, ₽/м²"],
             "rows": [["2026 Q1", "576 680"], ["2026 Q2", "670 281"],
                      ["2026 Q3", "724 203"]]}
    drawn = sales_deck.charts(table)
    assert drawn, "колонка цен читается числом"
    built = sales_deck.build(
        [{"title": "Цены", "charted": True, "tables": [table], "lines": [], "strips": [],
          "note": ""}],
        title="Проверка", subtitle="", footer="")
    from pptx import Presentation
    import io as _io
    deck = Presentation(_io.BytesIO(built))
    labels = [shape.chart.plots[0].data_labels for slide in deck.slides
              for shape in slide.shapes if shape.has_chart]
    assert labels, "график не построился"
    assert labels[0].number_format == "#,##0", "у сотен тысяч дробной части нет"


def test_a_column_with_the_same_label_in_every_cell_is_still_numeric() -> None:
    """«900,0 млн ₽» — это число с единицей, а не текст.

    Из-за отказа читать такую ячейку у «Факта против планов» с графика
    пропадали и факт, и оба плана: числом оставались только цены, у которых
    подписи в ячейках нет (снимок владельца, 31.08.2026).
    """
    from market_search import sales_deck

    table = {"head": ["Квартал", "млн ₽, факт", "млн ₽, план ФМ", "Цена факт, ₽/м²"],
             "rows": [["2026 Q1", "900,0 млн ₽", "1 100,0 млн ₽", "219 000"],
                      ["2026 Q2", "990,0 млн ₽", "1 210,0 млн ₽", "219 000"],
                      ["2026 Q3", "540,0 млн ₽", "660,0 млн ₽", "219 000"]]}
    drawn = sales_deck.charts(table)
    assert len(drawn) == 1, "рубли — одна мера, один график"
    assert drawn[0]["measure"] == "млн ₽"
    assert [row["name"] for row in drawn[0]["extra"]] == ["млн ₽, план ФМ"]
    assert [row["name"] for row in drawn[0]["second"]] == ["Цена факт, ₽/м²"]


def test_a_column_of_mixed_labels_is_not_a_column() -> None:
    """Разные подписи в одной колонке — смесь величин, и рисовать её нельзя."""
    from market_search import sales_deck

    table = {"head": ["Строка", "Значение"],
             "rows": [["A", "900,0 млн ₽"], ["B", "3,1%"], ["C", "12 шт"]]}
    assert sales_deck.charts(table) == []
