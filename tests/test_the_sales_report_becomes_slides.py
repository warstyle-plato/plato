"""Слайд — это раздел отчёта о продажах, и его можно править.

«Страниц PDF = слайд или раздел PDF = слайд» (владелец, 27.08.2026): у нас эти
два ответа совпали — раздел свода печатается со своей страницы.

Первая версия клала на слайд снимок раздела. Владелец, 29.08.2026: «этот отчёт
в редактируемом виде нужен отделу продаж, картинка никому не уперлась». Значит
на слайде настоящие объекты: заголовок, вывод текстом, таблица ячейками,
график с данными.

И при этом — никакой второй сборки «по тем же данным»: это была бы вторая
реализация отчёта о продажах, она разошлась бы с экраном молча, и обе выглядели
бы верными. Колода собирается из ТОЙ ЖЕ разметки, которой печатается PDF, и
каждое число слайда буквально взято со строки экрана.

Запуск: python3 -m pytest tests/test_the_sales_report_becomes_slides.py -q
"""

from __future__ import annotations

import io
import re
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tests"))

from market_search import sales_deck  # noqa: E402
from market_search.cabinet import cabinet_page  # noqa: E402

MARKUP = """
<div class="blockhead"><h2>Продажи проекта — Тестовый ЖК</h2>
<div class="noprint"><button id="salespdf">Скачать PDF</button></div></div>
<div class="salesnav"><a href="#sb-dyn">Динамика</a></div>
<div class="kv"><div><div class="muted">Договоров</div><div>76</div>
<div class="muted">за всё время</div></div>
<div><div class="muted">Выручка</div><div>1 628,9 млн ₽</div>
<div class="muted">12,4% из ожидаемых</div></div></div>
<div class="sumup">Продано 12,4% ожидаемой выручки проекта.</div>
<section id="sb-dyn" class="salesblock"><div class="blockhead"><h3>Динамика</h3>
<div class="switch"><button class="on">млн ₽</button><button>лоты</button></div></div>
<div class="wrap"><svg><text>не текст слайда</text></svg></div>
<div class="muted">млн ₽ по месяцам</div>
<details><summary>Помесячно числами</summary>
<table><thead><tr><th>Месяц</th><th>Лотов</th><th>млн ₽</th><th>₽/м²</th></tr></thead>
<tbody><tr><td>2026-07</td><td>5</td><td>120,4</td><td>540 000</td></tr>
<tr><td>2026-06</td><td>3</td><td>84,0</td><td>—</td></tr></tbody></table></details>
<div class="sumup">Темп последних месяцев ниже плана.</div></section>
<div class="muted">Источники: контрактация ЦФ — 2026-08-20.</div>
<div class="muted">Не прочитано — план банка не загружен.</div>
"""


def parsed() -> list[dict]:
    return sales_deck.sections(MARKUP)


def test_the_deck_builder_touches_no_number_of_its_own() -> None:
    """Первое «просто посчитать долю» здесь — это вторая экономика продаж."""
    body = (ROOT / "market_search" / "sales_deck.py").read_text(encoding="utf-8")
    for name in ("amount", "escrow", "units", "contracts", "summarise",
                 "conclusions", "by_channel", "revenue"):
        assert name not in body, f"сборщик колоды знает про «{name}» — это второй счёт"
    assert "contracting" not in body


def test_the_dependency_is_declared_where_the_image_reads_it() -> None:
    assert "python-pptx" in (ROOT / "requirements.txt").read_text(encoding="utf-8")


def test_the_sections_of_the_screen_become_the_pages_of_the_deck() -> None:
    got = parsed()
    titles = [item["title"] for item in got]
    assert titles[0].startswith("Продажи проекта")
    assert "Динамика" in titles
    assert titles[-1] == "На чём посчитано", "оговорка осталась бы на экране"
    tail = got[-1]
    assert any("Источники" in line for line in tail["lines"])
    assert any("Не прочитано" in line for line in tail["lines"])


def test_the_screen_tools_do_not_reach_the_slide() -> None:
    """Кнопки, переключатель меры и навигация по якорям — орудия экрана."""
    flat = " ".join(line for item in parsed() for line in item["lines"])
    for tool in ("Скачать PDF", "млн ₽лоты", "Динамика</a>"):
        assert tool not in flat
    head = parsed()[0]
    assert not any("Динамика" == line for line in head["lines"]), "якорь уехал на слайд"


def test_the_key_numbers_come_as_a_table_not_as_loose_lines() -> None:
    """Россыпью строк плашки наезжали друг на друга и читались как обрывки."""
    head = parsed()[0]
    assert head["tables"], "плашка ключевых чисел не собралась"
    table = head["tables"][0]
    assert table["head"] == ["Показатель", "Значение", "Пояснение"]
    assert ["Договоров", "76", "за всё время"] in table["rows"]
    assert head["note"].startswith("Продано")


def test_a_cell_is_a_number_whole_or_not_at_all() -> None:
    assert sales_deck.cell_number("3 306 021") == 3306021.0
    assert sales_deck.cell_number("120,4") == 120.4
    for text in ("—", "", "5,00%", "1 628,9 млн ₽", "не число"):
        assert sales_deck.cell_number(text) is None, text


def test_every_measure_gets_its_own_chart() -> None:
    """«Не проще для каждого графика свой слайд сделать?» — проще (владелец,
    29.08.2026): каждая мера показана, ни одна не спорит с соседней на общей
    оси, а лишний слайд в PowerPoint удаляют одним нажатием.

    Экран предлагает меры переключателем; в документе переключателя нет, и
    мера, до которой не переключились, в нём просто отсутствует — то же
    правило, что у свёрнутой таблицы. Схлопывать их в один график было
    ошибкой: «куча столбиков» приходила из разделов, где графика нет вовсе.

    А колонка, в которой осталось меньше двух чисел, графиком не становится:
    по одной точке линии нет, и прочерк рядом — пропуск, а не ноль.
    """
    table = parsed()[1]["tables"][0]
    drawn = sales_deck.charts(table)
    # Деньги первыми: управленцу нужны рубли, штуки при них справочны.
    assert [item["name"] for item in drawn] == ["млн ₽", "Лотов"]
    # «₽/м²» здесь одна цифра и прочерк — линии из этого не выйдет.
    assert all(not item["second"] for item in drawn)
    assert drawn[0]["categories"] == ["2026-07", "2026-06"]
    # Один ряд на график: рубли и штуки на одной оси дают столбик в пиксель.
    assert all(len(item["values"]) == len(item["categories"]) for item in drawn)


def test_the_slides_carry_real_objects_and_not_a_single_picture() -> None:
    pptx = pytest.importorskip("pptx")

    raw = sales_deck.build(parsed(), title="Продажи — Тестовый ЖК",
                           subtitle="срез 2026-08-20", footer="DevelopAid")
    deck = pptx.Presentation(io.BytesIO(raw))
    assert deck.slide_width > deck.slide_height, "слайд 16:9"

    tables, charts, pictures, texts = 0, 0, 0, []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                pictures += 1
            elif getattr(shape, "has_table", False) and shape.has_table:
                tables += 1
            elif getattr(shape, "has_chart", False) and shape.has_chart:
                charts += 1
            elif shape.has_text_frame:
                texts.append(shape.text_frame.text)
    assert pictures == 0, "картинка вместо объектов — её нельзя править"
    assert tables >= 1 and charts >= 1, "раздел с числами остался без графика"
    assert any("Продажи — Тестовый ЖК" in text for text in texts)
    assert any("Темп последних месяцев" in text for text in texts), "вывод раздела"

    # Числа на слайде — те же, что в разметке экрана. Ключевые лежат плитками
    # (фигура с текстом), помесячные — ячейками таблицы.
    numbers = [cell.text for slide in deck.slides for shape in slide.shapes
               if getattr(shape, "has_table", False) and shape.has_table
               for row in shape.table.rows for cell in row.cells]
    assert "120,4" in numbers and "540 000" in numbers
    assert any("76" in text for text in texts), "ключевое число пропало с плитки"


def test_nothing_to_show_is_said_out_loud() -> None:
    pytest.importorskip("pptx")
    with pytest.raises(sales_deck.DeckUnavailable):
        sales_deck.build([], title="Продажи", subtitle="", footer="")
    with pytest.raises(sales_deck.DeckUnavailable):
        sales_deck.sections("   ")


def test_the_button_sends_the_same_markup_as_the_pdf() -> None:
    """Две сборки одной колоды разойдутся — значит разметка одна и та же."""
    body = max(re.findall(r"<script[^>]*>(.*?)</script>", cabinet_page("sales"), re.S), key=len)
    assert 'id="salesppt"' in cabinet_page("sales")
    handler = body[body.index("$('#salesppt').onclick="):][:2200]
    assert "salesPrintHtml()" in handler, "презентация собиралась бы из другой разметки"
    assert "'/cabinet/sales.pptx'" in handler
    assert "Презентация не собралась" in handler
    assert "window.print()" not in handler


def test_the_route_refuses_an_empty_body_with_a_reason(tmp_path, monkeypatch) -> None:
    from fastapi import FastAPI
    from fastapi.testclient import TestClient

    from market_search.api import install

    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    monkeypatch.setenv("MARKET_CABINET_KEY", "stand-key-2026")
    app = FastAPI()
    install(app)
    client = TestClient(app)
    client.post("/cabinet/login", content="key=stand-key-2026",
                headers={"Content-Type": "application/x-www-form-urlencoded"},
                follow_redirects=False)

    answer = client.post("/cabinet/sales.pptx", json={"html": "   "})
    assert answer.status_code == 422, answer.text
    assert "пуст" in answer.json()["detail"]

    monkeypatch.delenv("MARKET_CABINET_KEY", raising=False)
    closed = TestClient(app).post("/cabinet/sales.pptx", json={"html": "<b>x</b>"})
    assert closed.status_code == 503 and "Кабинет" in closed.json()["detail"]


def test_the_live_screen_markup_parses_into_slides(tmp_path) -> None:
    """Разбор проверяется на настоящей разметке экрана, а не только на образце.

    Образец пишу я, и он сойдётся с разбором даже когда экран сменил вёрстку.
    Без Chromium — пропуск, а не зелёный прогон на пустом месте.
    """
    pytest.importorskip("pptx")
    play = pytest.importorskip("playwright.sync_api")
    import importlib

    import browser_launch

    from market_search import contracting

    got = importlib.import_module("test_contracting_summary")._summary()
    got["sources"] = [{"kind": "contracting", "name": "контрактация ЦФ",
                       "at": "2026-08-20T10:00:00"}]
    got["plans"] = contracting.plan_comparison(got)
    got["conclusions"] = contracting.conclusions(got)
    got["pool"] = contracting.pool_progress(got, [], None, None)

    file = tmp_path / "cabinet.html"
    file.write_text(cabinet_page("sales").replace("__DEVELOPAID_VERSION__", "test"), encoding="utf-8")
    with play.sync_playwright() as pw:
        try:
            browser = browser_launch.launch(pw)
        except Exception as exc:  # noqa: BLE001
            pytest.skip(f"Chromium недоступен: {exc}")
        try:
            tab = browser.new_page()
            tab.route("**/*", lambda route: route.abort()
                      if route.request.url.startswith("http") else route.continue_())
            tab.goto(file.as_uri())
            markup = tab.evaluate("(d)=>{renderSales(d); return salesPrintHtml()}", got)
        finally:
            browser.close()

    blocks = sales_deck.sections(markup)
    titles = [item["title"] for item in blocks]
    assert titles[-1] == "На чём посчитано"
    assert markup.count('class="salesblock"') + 2 == len(blocks), \
        "разделы экрана и листы колоды разошлись"
    assert blocks[0]["tables"], "ключевые числа не доехали таблицей"
    raw = sales_deck.build(blocks, title="Продажи", subtitle="срез", footer="DevelopAid")
    assert raw[:2] == b"PK", "это не .pptx"


# --- Лист без содержимого листом не является ---------------------------------
#
# Снимок колоды с телефона (владелец, 31.08.2026): слайд 2 — заголовок и одна
# строка вывода, слайд 3 — заголовок, подпись графика и плашка ПОВЕРХ неё, а
# сами графики следом отдельными листами. На экране эти подписи стоят рядом с
# картинкой, и отрывать их от неё незачем.

def _deck_slides(html: str):
    from pptx import Presentation

    from market_search import sales_deck

    raw = sales_deck.build(sales_deck.sections(html), title="Продажи — Проект",
                           subtitle="Свод продаж DevelopAid · срез 2026-08-31",
                           footer="")
    return Presentation(io.BytesIO(raw))


def _text_boxes(slide):
    from pptx.util import Emu

    out = []
    for shape in slide.shapes:
        text = shape.text_frame.text.strip() if shape.has_text_frame else ""
        if text:
            out.append((Emu(shape.left).inches, Emu(shape.top).inches,
                        Emu(shape.width).inches, Emu(shape.height).inches, text))
    return out


def test_a_chart_section_opens_with_its_chart() -> None:
    """Подпись графика — не слайд: она едет на слайд с графиком."""
    deck = _deck_slides(MARKUP)
    for index, slide in enumerate(deck.slides):
        if index == 0:
            continue
        heading = next((s.text_frame.text for s in slide.shapes
                        if s.has_text_frame and s.text_frame.text), "")
        if "Динамика" in heading:
            assert any(s.has_chart for s in slide.shapes), \
                "первый слайд «Динамики» обязан нести картинку, а не одну подпись"
            break
    else:
        pytest.skip("в этом своде нет раздела с графиком")


def test_no_text_sits_on_top_of_other_text() -> None:
    """Плашка вывода ложилась поверх подписи: `top` не двигался после подписей."""
    deck = _deck_slides(MARKUP)
    clashes = []
    for number, slide in enumerate(deck.slides, 1):
        boxes = _text_boxes(slide)
        for left in range(len(boxes)):
            for right in range(left + 1, len(boxes)):
                x1, y1, w1, h1, t1 = boxes[left]
                x2, y2, w2, h2, t2 = boxes[right]
                # Колонтитул и номер листа делят одну строку намеренно: левый
                # прижат влево, правый вправо.
                if abs(y1 - y2) < 0.01 and (t1.isdigit() or t2.isdigit()):
                    continue
                if (min(x1 + w1, x2 + w2) - max(x1, x2) > 0.05
                        and min(y1 + h1, y2 + h2) - max(y1, y2) > 0.05):
                    clashes.append((number, t1[:40], t2[:40]))
    assert not clashes, f"текст поверх текста: {clashes[:3]}"


def test_the_footer_says_each_thing_once() -> None:
    """Вызывающий слал имя и дату, сборщик добавлял их же — строка повторяла себя."""
    deck = _deck_slides(MARKUP)
    line = next(s.text_frame.text for s in list(deck.slides)[1].shapes
                if s.has_text_frame and "срез" in s.text_frame.text)
    parts = [p.strip() for p in line.split(" · ") if p.strip()]
    assert len(parts) == len(set(p.casefold() for p in parts)), f"повтор в колонтитуле: {line}"
    assert "слайд отвечает" not in line, "колонтитул объясняет устройство колоды"
