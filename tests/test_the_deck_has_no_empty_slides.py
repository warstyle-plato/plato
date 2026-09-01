"""Слайд не бывает пустым, а «Факт против планов» показывает планы.

«Пустые два слайда, кривой план-факт, хотя на сайте он прекрасен» (владелец,
31.08.2026). Проверено сборкой на настоящей вёрстке свода, и оба нашлись:

* раздел с тремя мерами давал три почти одинаковых столбиковых листа подряд —
  рубли, лоты, метры, — и на двух из них не было ничего, кроме картинки;
* таблица из двух строк ростом в 0,34 дюйма стояла под самой шапкой листа в
  7,5 — шесть дюймов белого читаются как «здесь ничего нет».

Рисуется ведущая мера, остальные названы под графиком и живут колонками
таблицы; короткий блок тянется и встаёт по центру, а таблица раздела садится
на лист графика, когда помещается.

И отдельно — ловушка чтения: у комбинированного графика ряды лежат в РАЗНЫХ
plot'ах, и `chart.plots[0].series` показывает один ряд из шести. Проверять надо
все планы разом, иначе полный график выглядит пустым (я на этом ошибся сам).

Запуск: python3 -m pytest tests/test_the_deck_has_no_empty_slides.py -q
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from market_search import sales_deck  # noqa: E402

SHEET_TOP, SHEET_BOTTOM = sales_deck.CONTENT_TOP, 6.7

PLANS = """
<section class="salesblock"><div class="blockhead"><h3>Факт против планов</h3></div>
<svg viewBox="0 0 900 300"><rect width="10" height="10"/></svg>
<details open><summary>Кварталы числами</summary><table><tbody>
<tr><th>Квартал</th><th>млн ₽, факт</th><th>млн ₽, план ФМ</th><th>млн ₽, план банка</th>
<th>Цена факт, ₽/м²</th><th>Цена ФМ, ₽/м²</th><th>Цена банка, ₽/м²</th></tr>
<tr><td>2025 Q3</td><td>420,0 млн ₽</td><td>500,0 млн ₽</td><td>610,0 млн ₽</td>
<td>268 000</td><td>275 000</td><td>281 000</td></tr>
<tr><td>2025 Q4</td><td>600,0 млн ₽</td><td>710,0 млн ₽</td><td>850,0 млн ₽</td>
<td>273 200</td><td>279 800</td><td>286 000</td></tr>
<tr><td>2026 Q1</td><td>780,0 млн ₽</td><td>920,0 млн ₽</td><td>1 090,0 млн ₽</td>
<td>278 400</td><td>284 600</td><td>291 000</td></tr>
</tbody></table></details>
<div class="sumup">Факт ниже обоих планов.</div></section>
<section class="salesblock"><div class="blockhead"><h3>Динамика</h3></div>
<svg viewBox="0 0 900 300"><rect width="10" height="10"/></svg>
<details open><summary>Помесячно числами</summary><table><tbody>
<tr><th>Месяц</th><th>Лотов</th><th>м²</th><th>млн ₽</th><th>₽/м²</th></tr>
<tr><td>2025-08</td><td>2</td><td>100</td><td>54,0</td><td>540 000</td></tr>
<tr><td>2025-09</td><td>1</td><td>14</td><td>4,0</td><td>285 714</td></tr>
</tbody></table></details></section>
<section class="salesblock"><div class="blockhead"><h3>Расторжения</h3></div>
<table><tbody><tr><th>Договор</th><th>Дата</th><th>млн ₽</th></tr>
<tr><td>1-2-9/ГР</td><td>2026-02-12</td><td>49,80</td></tr></tbody></table></section>
"""


def _deck():
    pytest.importorskip("pptx")
    from pptx import Presentation

    raw = sales_deck.build(sales_deck.sections(PLANS), title="Продажи — Проба",
                           subtitle="срез", footer="DevelopAid")
    return Presentation(io.BytesIO(raw))


def _band(slide) -> tuple[float, float]:
    """Верх и низ содержимого слайда — без шапки и колонтитула."""
    tops, bottoms = [], []
    for shape in slide.shapes:
        top = shape.top / 914400
        bottom = top + shape.height / 914400
        if top < 1.3 or top > SHEET_BOTTOM:
            continue
        tops.append(top)
        bottoms.append(bottom)
    return (min(tops), max(bottoms)) if tops else (0.0, 0.0)


def test_the_plan_chart_carries_both_plans():
    """Ряды комбинированного графика лежат в разных plot'ах — считаем все."""
    charted = []
    for slide in _deck().slides:
        for shape in slide.shapes:
            if shape.has_chart:
                charted.append([series.name for plot in shape.chart.plots
                                for series in plot.series])
    assert charted, "график не построен вовсе"
    plans = next((names for names in charted if any("факт" in n for n in names)), [])
    assert any("план ФМ" in name for name in plans), \
        "на слайде «Факт против планов» нет плана ФМ"
    assert any("план банка" in name for name in plans), \
        "на слайде «Факт против планов» нет плана банка"
    assert any("Цена" in name for name in plans), "цена не доехала на вторую шкалу"


def test_one_section_does_not_become_three_alike_slides():
    titles = []
    for slide in _deck().slides:
        for shape in slide.shapes:
            if shape.has_text_frame and 0.8 < shape.top / 914400 < 1.0:
                titles.append(shape.text_frame.paragraphs[0].text)
                break
    charts = [name for name in titles if "Динамика ·" in name]
    assert len(charts) <= 1, f"одна динамика разъехалась на листы: {charts}"


def test_no_slide_is_a_strip_glued_under_the_heading():
    """Короткий блок тянется и встаёт по центру — иначе это пустой лист."""
    room = SHEET_BOTTOM - SHEET_TOP
    for index, slide in enumerate(_deck().slides, 1):
        top, bottom = _band(slide)
        if not top:
            continue
        filled = (bottom - top) / room
        if filled >= 0.35:
            continue
        # Меньше трети листа — тогда блок обязан стоять по центру, а не
        # прижатым к шапке: шесть дюймов белого под строкой читаются как
        # «здесь ничего нет».
        above, below = top - SHEET_TOP, SHEET_BOTTOM - bottom
        assert above > 0.6 and abs(above - below) < 1.2, (
            f"слайд {index}: блок {filled:.0%} листа прижат к шапке "
            f"(сверху {above:.2f}\", снизу {below:.2f}\")")


def test_nothing_runs_off_the_sheet_and_nothing_overlaps():
    for index, slide in enumerate(_deck().slides, 1):
        boxes = []
        for shape in slide.shapes:
            top = shape.top / 914400
            bottom = top + shape.height / 914400
            if shape.has_text_frame and not shape.text_frame.text.strip():
                continue
            assert bottom <= 7.4, f"слайд {index}: объект ушёл за нижний край"
            boxes.append((top, bottom, shape.left / 914400,
                          (shape.left + shape.width) / 914400))
        for first in range(len(boxes)):
            for second in range(first + 1, len(boxes)):
                t1, b1, l1, r1 = boxes[first]
                t2, b2, l2, r2 = boxes[second]
                assert not (min(b1, b2) - max(t1, t2) > 0.08
                            and min(r1, r2) - max(l1, l2) > 0.3), \
                    f"слайд {index}: объекты легли друг на друга"
