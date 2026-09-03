"""Воронка обращений: верх, которого в своде не было вовсе.

Свод продаж начинался с подписанного договора, то есть с конца. Обращения
объясняют остальное: девять из десяти — звонки с конверсией в бронь около трёх
процентов, тогда как агентские каналы дают тридцать шесть.

«Эти данные надо дать Платону к выводам» (владелец, 26.08.2026) — поэтому
воронка идёт и в вывод под блоком, и в вопрос Платону.
"""
from __future__ import annotations

from pathlib import Path

from market_search import contracting, demand

CABINET = Path(__file__).resolve().parent.parent / "market_search" / "cabinet.py"


def _deals(rows: list[dict]) -> list[dict]:
    base = {"source": "Звонок", "manager": "Иванов", "booked": False,
            "need_asked": False, "next_step": False, "not_a_lead": False,
            "area_min": None, "budget_max": None, "wants": []}
    return [{**base, **row} for row in rows]


def test_a_non_lead_call_is_out_of_the_denominator() -> None:
    """Реклама и предложения услуг — не обращение, и доля по ним врёт."""
    got = demand.funnel(_deals([
        {"not_a_lead": True}, {"booked": True}, {},
    ]))
    assert got["quality"]["calls"] == 3
    assert got["quality"]["not_a_lead"] == 1
    assert got["quality"]["target"] == 2
    assert got["quality"]["booked_target"] == 0.5


def test_the_card_is_measured_not_the_call() -> None:
    """Комментарий — пересказ BitrixGPT, и это сказано вслух."""
    got = demand.funnel([])
    assert any("пересказ BitrixGPT" in note for note in got["notes"])
    assert any("перезвонить не с чем" in note for note in got["notes"])


def test_what_the_funnel_cannot_do_is_said() -> None:
    """Связи с договором нет, денег у обращения нет — обещать их нельзя."""
    got = demand.funnel([])
    assert any("ни номера договора" in note for note in got["notes"])
    assert any("сколько выручки" in note for note in got["notes"])
    assert any("ещё рано" in note for note in got["notes"])


def test_the_share_travels_with_its_count() -> None:
    """На пяти бронях доля не значит ничего, и число обращений рядом."""
    got = demand.funnel(_deals([{"source": "Сайт"}, {"source": "Сайт", "booked": True}]))
    row = next(x for x in got["by_source"] if x["name"] == "Сайт")
    assert row["deals"] == 2 and row["booked"] == 1 and row["share"] == 0.5


def test_the_conclusion_calls_the_gap_a_neighbourhood_not_a_cause() -> None:
    """Менеджер мог расспрашивать тех, кто и так был готов."""
    deals = _deals(
        [{"need_asked": True, "booked": True}] * 3
        + [{"need_asked": True}] * 7
        + [{}] * 40)
    summary = {"demand": {"funnel": demand.funnel(deals)}, "total": {}, "dynamics": []}
    line = contracting.conclusions(summary)["funnel"]
    assert "не доказанная причина" in line
    assert "не осталось ни потребности" in line


def test_the_funnel_goes_into_the_question() -> None:
    page = CABINET.read_text()
    body = page[page.index("function salesDigest("):page.index("async function askPlatoSales(")]
    assert "ВОРОНКА:" in body and "ИСТОЧНИК ${x.name}" in body
    assert "МЕНЕДЖЕР ${x.name}" in body


def test_a_block_without_a_conclusion_says_what_is_missing() -> None:
    """Пустое место под блоком и отсутствующий вывод выглядят одинаково."""
    page = CABINET.read_text()
    assert "const NOTE_NEEDS=" in page
    body = page[page.index("function salesNote("):page.index("\n}", page.index("function salesNote("))]
    assert "Вывод не сложился" in body
    for key in ("pool", "bands", "demand", "funnel", "escrow", "fm", "bank"):
        assert f"{key}:" in page[page.index("const NOTE_NEEDS="):page.index("function salesNote(")], key


def test_the_funnel_is_drawn_and_the_picture_reaches_the_deck(tmp_path) -> None:
    """«У тебя таблица воронки обращений идёт!!!» (владелец, 03.09.2026).

    Таблица шла, картинки не было: раздел рисовался одними строками, и в
    презентации от воронки оставались колонки — разрыв между звонками и
    агентами приходилось искать глазами. Колода рисует ровно там, где рисует
    экран, поэтому чинится это на экране, а на слайд приходит само.

    Проверяется на живой странице, а не поиском строки: график собирает общий
    `barChart`, и обещание «здесь картинка» стоит ровно столько, сколько
    стоит его вызов на настоящих числах.
    """
    import importlib
    import sys

    import pytest

    play = pytest.importorskip("playwright.sync_api")
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    import browser_launch

    from market_search import cabinet, sales_deck

    got = importlib.import_module(
        "test_the_sales_report_survives_a_full_project").full_summary()
    lead = got.setdefault("demand", {}).setdefault("funnel", {})
    lead["by_source"] = [{"name": "Звонок", "deals": 518, "booked": 16, "share": 0.031},
                         {"name": "Агент", "deals": 44, "booked": 16, "share": 0.364},
                         {"name": "Сайт", "deals": 11, "booked": 1, "share": 0.091}]
    lead["by_manager"] = [{"name": "Иванова", "deals": 200, "booked": 25, "share": 0.125},
                          {"name": "Петров", "deals": 160, "booked": 2, "share": 0.012}]
    lead.setdefault("quality", {"calls": 573, "target": 449,
                                "booked_target": 0.031, "blank": 231})

    file = tmp_path / "cabinet.html"
    file.write_text(cabinet.cabinet_page("sales").replace("__DEVELOPAID_VERSION__", "t"),
                    encoding="utf-8")
    with play.sync_playwright() as pw:
        try:
            browser = browser_launch.launch(pw)
        except Exception as exc:  # noqa: BLE001
            pytest.skip(f"Chromium недоступен: {exc}")
        try:
            tab = browser.new_page()
            tab.route("**/*", lambda route: route.abort()
                      if route.request.url.startswith("http") else route.continue_())
            tab.goto(file.as_uri())
            markup = tab.evaluate("(d)=>{renderSales(d); return salesPrintHtml()}", got)
        finally:
            browser.close()

    block = next(page for page in sales_deck.sections(markup)
                 if str(page.get("title") or "").startswith("Воронка"))
    drawn = [table for table in (block.get("tables") or []) if table.get("charted")]
    assert [table["head"][0] for table in drawn] == ["Источник", "Менеджер"], \
        "картинки нет ни у источников, ни у менеджеров"
    for table in drawn:
        charts = sales_deck.charts(table)
        assert len(charts) == 1, "воронка — одна картинка, а не лист на колонку"
        rows = [charts[0]["name"], *[extra["name"] for extra in charts[0]["extra"]]]
        assert rows == ["Обращений", "Броней"], rows
        # Доля — линия справа на своей шкале, как на экране: своим листом со
        # столбиками она читалась бы как ещё один объём.
        assert [line["name"] for line in charts[0]["second"]] == ["Доля"]

    import io

    from pptx import Presentation

    deck = Presentation(io.BytesIO(sales_deck.build(
        sales_deck.sections(markup), title="Т", subtitle="с", footer="ф")))
    headings = [shape.text_frame.text for slide in deck.slides for shape in slide.shapes
                if shape.has_text_frame and 700000 < shape.top < 1_150_000]
    funnel = [line for line in headings if line.startswith("Воронка обращений ·")
              and "продолжение" not in line]
    assert len(funnel) == 2 and len(set(funnel)) == 2, \
        f"два листа воронки нечем различить: {funnel}"
