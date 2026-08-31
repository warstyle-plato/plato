"""Предпросмотр колоды картинками: фигуры по их координатам, графики по данным.

LibreOffice в песочнице pptx не открывает, `pdftoppm` нет — увидеть слайд
иначе нечем, а «убого или нет» разбором файла не проверяется. Это не рендер
PowerPoint: пустоту, переполнение, выравнивание и общий вид показывает, шрифты
и графики приблизительны.
"""
from __future__ import annotations

import html
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

PX = 96.0  # пикселей на дюйм


def inches(value) -> float:
    return Emu(value).inches if value is not None else 0.0


def colour(fmt) -> str:
    try:
        return "#%02X%02X%02X" % tuple(fmt.rgb)
    except Exception:  # noqa: BLE001
        return "#16202B"


def chart_svg(chart, width: float, height: float) -> str:
    """Рисуем ВСЕ группы: столбики, линии своей шкалы и линии второй.

    Пока рисовалась первая группа, сводный график «факт против планов»
    выглядел одним фактом — и я чинил бы то, что в файле уже верно.
    """
    plots = list(chart.plots)
    cats = [str(c) for c in plots[0].categories]
    bars = [(s.name, list(s.values)) for s in plots[0].series]
    own = [a.get("val") for a in plots[0]._element.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/chart}axId")]
    left, right = [], []
    for plot in plots[1:]:
        axes = [a.get("val") for a in plot._element.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}axId")]
        rows = [(s.name, list(s.values)) for s in plot.series]
        (left if axes == own else right).extend(rows)
    tones = ["#1367AE", "#7FB2E5", "#B9CFE4", "#D7E4F0"]
    left_tones = ["#C4581B", "#7C6BB5", "#8A9BA8"]
    right_tones = ["#1F5C87", "#D9A441", "#4FA07A", "#9A6BB5"]
    pad = 40
    w, h = width - pad * 2, height - pad * 2 - (26 if (left or right) else 0)
    def span(rows):
        vals = [v for _, values in rows for v in values if v is not None]
        return (min(vals), max(vals)) if vals else (0.0, 1.0)
    top = max([v for _, values in bars + left for v in values if v is not None] or [1]) or 1
    rlo, rhi = span(right)
    rlo = rlo * 0.9
    step = w / max(1, len(cats))
    slot = step / (1 + (plots[0].gap_width or 150) / 100.0)
    bar_w = slot / max(1, len(bars))
    out = [f'<svg width="{width}" height="{height}">']
    def x_at(i): return pad + i * step + step / 2
    for i in range(len(cats)):
        base = pad + i * step + (step - slot) / 2
        for order, (_, values) in enumerate(bars):
            value = values[i] if i < len(values) else None
            if value is None:
                continue
            bh = h * (value / top)
            x = base + order * bar_w
            out.append(f'<rect x="{x:.1f}" y="{pad + h - bh:.1f}" width="{bar_w:.1f}"'
                       f' height="{bh:.1f}" fill="{tones[min(order, 3)]}"/>')
            if plots[0].has_data_labels:
                out.append(f'<text x="{x + bar_w / 2:.1f}" y="{pad + h - bh - 6:.1f}"'
                           f' font-size="11" font-weight="700" text-anchor="middle"'
                           f' fill="#16202B">{value:,.0f}</text>'.replace(",", " "))
        out.append(f'<text x="{base + slot / 2:.1f}" y="{pad + h + 16:.1f}" font-size="10"'
                   f' text-anchor="middle" fill="#5B6B7D">{html.escape(cats[i])}</text>')
    def polyline(rows, palette, lo, hi):
        for order, (_, values) in enumerate(rows):
            pts = [f"{x_at(i):.1f},{pad + h - h * 0.88 * ((v - lo) / ((hi - lo) or 1)) - h * 0.06:.1f}"
                   for i, v in enumerate(values) if v is not None]
            if len(pts) < 2:
                continue
            colour = palette[min(order, len(palette) - 1)]
            out.append(f'<polyline points="{" ".join(pts)}" fill="none"'
                       f' stroke="{colour}" stroke-width="2.5"/>')
            for point in pts:
                cx, cy = point.split(",")
                out.append(f'<circle cx="{cx}" cy="{cy}" r="3" fill="{colour}"/>')
    polyline(left, left_tones, 0.0, top)
    polyline(right, right_tones, rlo, rhi)
    if chart.has_legend:
        marks = []
        every = bars + [(n, v) for n, v in left] + [(n, v) for n, v in right]
        palette = tones[:len(bars)] + left_tones[:len(left)] + right_tones[:len(right)]
        at = pad
        for order, (name, _) in enumerate(every):
            marks.append(f'<rect x="{at}" y="{height - 22}" width="10" height="10"'
                         f' fill="{palette[min(order, len(palette) - 1)]}"/>'
                         f'<text x="{at + 14}" y="{height - 13}" font-size="10"'
                         f' fill="#5B6B7D">{html.escape(str(name))}</text>')
            at += 34 + 6.2 * len(str(name))
        out.extend(marks)
    out.append("</svg>")
    return "".join(out)


def render(path: str, out_dir: str) -> list[str]:
    deck = Presentation(path)
    width, height = inches(deck.slide_width) * PX, inches(deck.slide_height) * PX
    made = []
    for number, slide in enumerate(deck.slides, 1):
        back = "#FFFFFF"
        try:
            if slide.background.fill.type is not None:
                back = colour(slide.background.fill.fore_color)
        except Exception:  # noqa: BLE001
            pass
        # Картинки рисуются: без них титул выглядит так же, как титул без
        # эмблемы, и починка не видна.
        parts = [f'<div class="slide" style="width:{width}px;height:{height}px;'
                 f'background:{back}">']
        for shape in slide.shapes:
            left, top = inches(shape.left) * PX, inches(shape.top) * PX
            w, h = inches(shape.width) * PX, inches(shape.height) * PX
            style = f"position:absolute;left:{left:.1f}px;top:{top:.1f}px;width:{w:.1f}px"
            if shape.__class__.__name__ == "Picture":
                import base64 as _b64
                blob = shape.image.blob
                kind = shape.image.content_type or "image/png"
                data = _b64.b64encode(blob).decode("ascii")
                parts.append(
                    f'<img src="data:{kind};base64,{data}" style="position:absolute;'
                    f'left:{inches(shape.left) * PX}px;top:{inches(shape.top) * PX}px;'
                    f'width:{inches(shape.width) * PX}px;height:{inches(shape.height) * PX}px">')
                continue
            if shape.has_chart:
                parts.append(f'<div style="{style};height:{h:.1f}px">'
                             + chart_svg(shape.chart, w, h) + "</div>")
            elif shape.has_table:
                rows = []
                for line_no, row in enumerate(shape.table.rows):
                    cells = []
                    for cell in row.cells:
                        fill = colour(cell.fill.fore_color) if cell.fill.type else "#FFF"
                        bold = "font-weight:700;" if line_no == 0 else ""
                        align = ("text-align:right;"
                                 if cell.text_frame.paragraphs[0].alignment else "")
                        cells.append(f'<td style="background:{fill};{bold}{align}">'
                                     f"{html.escape(cell.text)}</td>")
                    rows.append("<tr>" + "".join(cells) + "</tr>")
                parts.append(f'<table style="{style}">' + "".join(rows) + "</table>")
            elif str(shape.shape_type or "").startswith("AUTO_SHAPE"):
                try:
                    fill = colour(shape.fill.fore_color) if shape.fill.type is not None else "#FFF"
                except Exception:  # noqa: BLE001
                    fill = "#FFF"
                inner = ""
                if shape.has_text_frame and shape.text_frame.text.strip():
                    rows = []
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            size = run.font.size.pt if run.font.size else 14
                            weight = 700 if run.font.bold else 400
                            tone = colour(run.font.color) if run.font.color else "#16202B"
                            rows.append(f'<div style="font-size:{size * 1.333:.1f}px;'
                                        f'font-weight:{weight};color:{tone}">'
                                        f"{html.escape(run.text)}</div>")
                    inner = "".join(rows)
                parts.append(f'<div style="{style};height:{h:.1f}px;background:{fill};'
                             f'border:1px solid #DDE5ED;box-sizing:border-box;'
                             f'padding:{"8px 12px" if inner else "0"}">{inner}</div>')
            elif shape.has_text_frame and shape.text_frame.text.strip():
                lines = []
                for para in shape.text_frame.paragraphs:
                    runs = []
                    for run in para.runs:
                        size = run.font.size.pt if run.font.size else 14
                        weight = 700 if run.font.bold else 400
                        tone = colour(run.font.color) if run.font.color else "#16202B"
                        runs.append(f'<span style="font-size:{size * 1.333:.1f}px;'
                                    f'font-weight:{weight};color:{tone}">'
                                    f"{html.escape(run.text)}</span>")
                    align = {1: "left", 2: "center", 3: "right"}.get(
                        int(para.alignment) if para.alignment is not None else 1, "left")
                    lines.append(f'<div style="text-align:{align}">' + "".join(runs) + "</div>")
                skin = ""
                try:
                    if shape.fill.type is not None:
                        skin = (f"background:{colour(shape.fill.fore_color)};"
                                f"border:1px solid #DDE5ED;padding:8px 12px;"
                                f"box-sizing:border-box;height:{h:.1f}px")
                except Exception:
                    pass
                parts.append(f'<div style="{style};{skin}">' + "".join(lines) + "</div>")
        parts.append(f'<div class="num">слайд {number}</div></div>')
        made.append("".join(parts))
    page = ("<style>body{margin:0;background:#8894a2;font-family:Calibri,Arial,sans-serif}"
            ".slide{position:relative;margin:0 auto 18px;overflow:hidden;box-sizing:border-box}"
            "table{border-collapse:collapse;font-size:14px}"
            "td{border:1px solid #DDE5ED;padding:6px 8px}"
            ".num{position:absolute;left:6px;bottom:2px;font-size:10px;color:#c8d2dc}"
            "</style>" + "".join(made))
    Path(out_dir, "preview.html").write_text(page, encoding="utf-8")
    return made


if __name__ == "__main__":
    where = sys.argv[2] if len(sys.argv) > 2 else "."
    made = render(sys.argv[1], where)
    print(f"слайдов: {len(made)} → {Path(where, 'preview.html')}")
