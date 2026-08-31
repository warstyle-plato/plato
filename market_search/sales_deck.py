"""Свод продаж презентацией: слайд — раздел отчёта, и его можно править.

«Страниц PDF = слайд или раздел PDF = слайд» (владелец, 27.08.2026). Разделы и
страницы у нас совпадают — раздел свода печатается со своей страницы, — поэтому
вопрос решён одним ответом: слайд отвечает разделу.

Первая версия клала на слайд СНИМОК раздела. Владелец: «этот отчёт в
редактируемом виде нужен отделу продаж, картинка никому не уперлась»
(29.08.2026). Картинку нельзя ни поправить перед встречей, ни пересобрать под
свой шаблон — а именно это с колодой и делают.

Значит на слайде настоящие объекты PowerPoint: заголовок, вывод текстом,
таблица ячейками, график с данными. И вот чего при этом делать нельзя:
собирать их «по тем же данным» из свода. Это была бы вторая реализация отчёта о
продажах — она разошлась бы с экраном молча, и обе выглядели бы верными (так
уже расходились бот с сайтом, отчёт с книгой и книга с движком).

Поэтому колода собирается из ТОЙ ЖЕ разметки, которой печатается PDF: каждое
число слайда буквально взято из строки экрана. Считать здесь нечего и нечем —
в модуле нет ни одного имени величины свода, и это проверяется тестом. График
строится из своей же таблицы, а не считается заново: разойтись им негде.
"""

from __future__ import annotations

import io
import re
from html.parser import HTMLParser
from typing import Any

# Слайд 16:9. Дюймы, потому что в них считает сам формат.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
# Строк таблицы на слайде. Тринадцать месяцев в одну таблицу влезают, дальше
# таблица продолжается на следующем слайде: ужать её до нечитаемого — не то же
# самое, что поместить.
ROWS_PER_SLIDE = 13
# Текстовых строк раздела на слайде. Раздел, у которого таблицы нет вовсе
# (полосы долей), живёт своими подписями — они и есть его числа.
LINES_PER_SLIDE = 10
# Содержимое начинается под шапкой — надзаголовок, линейка, заголовок.
CONTENT_TOP = 1.5


class DeckUnavailable(RuntimeError):
    """Колода не собралась. Это ответ, а не повод отдать пустой файл."""


class _Sections(HTMLParser):
    """Разбор печатной разметки свода на разделы, таблицы и строки.

    Разметка своя, поэтому разбираем её сами: чужой библиотеки ради десятка
    тегов в образ не возим. Промах здесь не «пустой слайд», а слайд без
    половины чисел, поэтому пустой разбор объявляется отказом выше.
    """

    # `summary` — подпись сворачивалки, а не содержание: на слайде «Помесячно
    # числами» это осиротевшая фраза, под которой ничего нет.
    _SKIP = {"script", "style", "button", "textarea", "select", "svg", "summary"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.sections: list[dict[str, Any]] = []
        self.head: dict[str, Any] = {"title": "", "lines": [], "tables": [], "strips": [],
                                     "note": "", "charted": False}
        self.tail: dict[str, Any] = {"title": "На чём посчитано", "lines": [],
                                     "tables": [], "strips": [], "note": "",
                                     "charted": False}
        self._current = self.head
        self._kv_depth = 0
        self._kv_row: list[str] | None = None
        self._skip_depth = 0
        self._svg_depth = 0
        self._cell: list[str] | None = None
        self._row: list[str] | None = None
        self._table: dict[str, Any] | None = None
        self._text: list[str] = []
        self._in_head_cell = False
        self._want: str = ""
        # Подписи легенды и полос лежат соседними `span` без единого пробела
        # между ними: «факт, млн ₽» и «цена квартир, ₽/м²» на слайде выходили
        # одним словом. Разделитель ставится ТОЛЬКО между соседями одного
        # уровня — внутри строки `span` разбивать нечего.
        self._span_depth = 0
        self._closed_span_at: int | None = None
        self._strip: dict[str, Any] | None = None
        self._strip_depth = 0
        # Легенда, стоящая под лентой: цветной квадратик и подпись за ним.
        # Куски ленты каналов подписи не несут вовсе, и связывает их с именами
        # только цвет — как и для читателя.
        self._legend: dict[str, Any] | None = None

    # --- служебное -----------------------------------------------------
    def _flush_line(self) -> None:
        line = re.sub(r"\s+", " ", "".join(self._text)).strip()
        self._text = []
        if not line:
            return
        if self._want == "title" and not self._current["title"]:
            self._current["title"] = line
            return
        if self._want == "note":
            self._current["note"] = line
            return
        self._current["lines"].append(line)

    # --- разбор --------------------------------------------------------
    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        classes = dict(attrs).get("class") or ""
        if tag == "svg":
            # Экран рисует график там, где график отвечает на вопрос раздела,
            # и таблицу там, где отвечает таблица. Сам рисунок со страницы не
            # переносится — переносится РЕШЕНИЕ: где на слайде быть графику.
            self._svg_depth += 1
            self._current["charted"] = True
            return
        if self._svg_depth:
            return
        # `printviews` — те же меры, что колода строит сама из таблицы раздела.
        # Взятые ещё раз, они задвоили бы легенды строками раздела; на бумаге
        # они нужны, потому что там переключателя нет, а здесь есть лист на
        # меру.
        if (tag in self._SKIP or "noprint" in classes or "salesnav" in classes
                or "printviews" in classes):
            self._skip_depth += 1
            return
        if self._skip_depth:
            # Вложенное внутри пропускаемого тоже считается: закрытие каждого
            # `div` глубину уменьшает, и без пары на открытии первый же
            # вложенный `</div>` снимал пропуск досрочно — хвост блока
            # доезжал строками раздела. Касается и `noprint`.
            if tag == "div" or tag in self._SKIP:
                self._skip_depth += 1
            return
        if self._kv_depth:
            self._kv_depth += 1
            if self._kv_depth == 2:
                self._flush_line()
                self._kv_row = []
            elif self._kv_depth == 3:
                self._flush_line()
            return
        if tag == "div" and "kv" in classes.split():
            # Плашка ключевых чисел: имя, число и пояснение под ним. На экране
            # это плитки с крупным числом; таблицей «Показатель / Значение /
            # Пояснение» они перестают быть плитками и читаются как список.
            self._flush_line()
            self._kv_depth = 1
            self._table = {"head": ["Показатель", "Значение", "Пояснение"],
                           "rows": [], "kind": "tiles"}
            return
        style = dict(attrs).get("style") or ""
        if tag == "div" and "display:flex" in style.replace(" ", "") and "height:" in style:
            self._flush_line()
            # Подпись ленты — строка прямо над ней: «Пул проекта · как
            # построено». Её уже прочитали, поэтому берём последнюю.
            caption = self._current["lines"].pop() if self._current["lines"] else ""
            self._strip = {"caption": caption, "parts": []}
            self._strip_depth = 1
            self._legend = None
            return
        if self._strip is not None:
            self._strip_depth += 1
            share = _SHARE.search(style)
            colour = _COLOUR.search(style)
            title = dict(attrs).get("title") or ""
            # Кусок берётся по ШИРИНЕ, а не по подписи. У лент каналов подписи
            # нет вовсе: имена стоят в легенде под лентой и связаны с кусками
            # только цветом — и лента опознавалась пустой, то есть раздел
            # «Каналы продаж» уходил на слайд без единой картинки. Имя
            # подбирается из легенды по цвету, ровно так же, как это читает
            # человек.
            if share:
                self._strip["parts"].append({
                    "name": title, "share": float(share.group(1)),
                    "colour": (colour.group(1) if colour else "4E9BDE").upper()})
            return
        if tag in {"table", "section", "h1", "h2", "h3", "h4"}:
            self._legend = None
        if tag == "section" and "salesblock" in classes:
            self._flush_line()
            self._current = {"title": "", "lines": [], "tables": [], "strips": [],
                             "note": "", "charted": False}
            self.sections.append(self._current)
            return
        if tag in {"h1", "h2", "h3", "h4"}:
            self._flush_line()
            self._want = "title"
            return
        if tag == "div" and "sumup" in classes:
            self._flush_line()
            self._want = "note"
            return
        if tag == "table":
            self._flush_line()
            self._table = {"head": [], "rows": []}
            return
        if tag == "tr" and self._table is not None:
            self._row = []
            return
        if tag in {"td", "th"} and self._table is not None:
            self._cell = []
            self._in_head_cell = tag == "th"
            return
        if tag == "span" and self._legend is not None:
            colour = _COLOUR.search(style)
            if colour and "width" in style.replace(" ", ""):
                self._legend["colour"] = colour.group(1).upper()
                return
        if tag == "span":
            if self._closed_span_at == self._span_depth and self._text:
                tail = "".join(self._text).rstrip()
                if tail and not tail.endswith(("·", ",", ";", ":")):
                    self._text.append(" · ")
            self._span_depth += 1
            self._closed_span_at = None
            return
        self._closed_span_at = None
        if tag in {"br", "p", "div", "li"}:
            self._flush_line()

    def handle_endtag(self, tag: str) -> None:
        if self._strip is not None and tag == "div":
            self._strip_depth -= 1
            if self._strip_depth == 0:
                if self._strip["parts"]:
                    self._current.setdefault("strips", []).append(self._strip)
                    self._legend = {"strip": self._strip, "colour": ""}
                self._strip = None
            return
        if tag == "svg":
            self._svg_depth = max(0, self._svg_depth - 1)
            return
        if self._svg_depth:
            return
        if self._skip_depth:
            if tag in self._SKIP or tag == "div":
                self._skip_depth = max(0, self._skip_depth - 1)
            return
        if tag in {"td", "th"} and self._cell is not None:
            text = re.sub(r"\s+", " ", "".join(self._cell)).strip()
            if self._row is not None:
                self._row.append(text)
            self._cell = None
            return
        if tag == "tr" and self._table is not None and self._row is not None:
            if self._in_head_cell and not self._table["head"]:
                self._table["head"] = self._row
            elif any(cell for cell in self._row):
                self._table["rows"].append(self._row)
            self._row = None
            self._in_head_cell = False
            return
        if tag == "table" and self._table is not None:
            if self._table["rows"]:
                self._current["tables"].append(self._table)
            self._table = None
            return
        if self._kv_depth:
            if self._kv_depth == 3 and self._kv_row is not None:
                line = re.sub(r"\s+", " ", "".join(self._text)).strip()
                self._text = []
                self._kv_row.append(line)
            elif self._kv_depth == 2 and self._kv_row is not None:
                if any(cell for cell in self._kv_row) and self._table is not None:
                    self._table["rows"].append(self._kv_row[:3])
                self._kv_row = None
            self._kv_depth -= 1
            if self._kv_depth == 0 and self._table is not None:
                if self._table["rows"]:
                    self._current["tables"].append(self._table)
                self._table = None
            return
        if tag == "section":
            self._flush_line()
            self._current = self.tail
            self._want = ""
            return
        if tag == "span":
            self._span_depth = max(0, self._span_depth - 1)
            self._closed_span_at = self._span_depth
            return
        self._closed_span_at = None
        if tag in {"h1", "h2", "h3", "h4", "p", "div", "li"}:
            self._flush_line()
            self._want = ""

    def _take_legend(self, text: str) -> bool:
        """Подпись легенды — имя куска ленты, а не строка раздела."""
        if self._legend is None or not self._legend["colour"]:
            return False
        colour = self._legend["colour"]
        self._legend["colour"] = ""
        named = False
        for part in self._legend["strip"]["parts"]:
            if part["colour"] == colour and not part["name"]:
                part["name"] = text
                named = True
        return named

    def handle_data(self, data: str) -> None:
        if self._skip_depth or self._svg_depth:
            return
        if self._cell is not None:
            self._cell.append(data)
            return
        text = re.sub(r"\s+", " ", data).strip()
        if text and self._take_legend(text):
            return
        self._text.append(data)

    def close(self) -> None:  # noqa: D102
        super().close()
        self._flush_line()


def sections(html: str) -> list[dict[str, Any]]:
    """Разделы печатной разметки: заголовок, вывод, строки и таблицы."""
    parser = _Sections()
    parser.feed(html)
    parser.close()
    # Лист «На чём посчитано» существует, только когда на нём что-то есть:
    # заголовок у него наш, и пустой он выдал бы пустую колоду за собранную.
    tail = [parser.tail] if (parser.tail["lines"] or parser.tail["tables"]) else []
    out = [parser.head] + parser.sections + tail
    kept = [item for item in out
            if item["lines"] or item["tables"] or item["strips"] or item["note"]]
    if not kept:
        raise DeckUnavailable("В разметке свода не нашлось ни одного раздела")
    return kept


_SHARE = re.compile(r"width:\s*([\d.]+)%")
_COLOUR = re.compile(r"background:\s*#([0-9a-fA-F]{6})")
_NUMBER = re.compile(r"^-?\d[\d  ]*(?:[.,]\d+)?$")
# Заголовок колонки цены метра: кроме «₽/м²» встречается «руб/м²» и «цена, ₽/м²».
_PRICE = re.compile(r"(₽|руб)\s*/\s*м", re.IGNORECASE)
# Денежная колонка: рубли — то, на что смотрят, а штуки и метры при них
# справочны. «₽/м²» сюда не попадает — её ловит `_PRICE` раньше.
_MONEY = re.compile(r"(₽|руб|млн|млрд)", re.IGNORECASE)


def cell_number(text: str) -> float | None:
    """Число из ячейки таблицы — целиком или никак.

    «3 306 021» уже читалось с середины в ценах рынка и во второй раз в
    комментариях CRM; здесь ячейка либо число целиком, либо не число.
    Проценты и суммы с подписями числом не считаются: график по колонке
    «5,00%» и по колонке «млн ₽» — разные графики.
    """
    raw = str(text or "").strip().replace("−", "-")
    if not raw or not _NUMBER.match(raw):
        return None
    try:
        return float(raw.replace(" ", "").replace(" ", "").replace(",", "."))
    except ValueError:
        return None


# Число с подписью: «3,1%», «1 628,9 млн ₽», «488 300 ₽/м²». Для ГРАФИКА это
# не число — колонка процентов и колонка рублей дают разные графики, и
# `cell_number` их не берёт намеренно. Но для ВЫКЛАДКИ это число: колонка
# «3,1% / 36,4% / 6,6%», прижатая влево, столбиком не читается, а её заголовок
# уже стоит справа. Два разных вопроса — две разные проверки.
_NUMERIC_CELL = re.compile(
    r"^-?\d[\d  ]*(?:[.,]\d+)?\s*(?:%|×|x|руб|₽|млн|млрд|тыс|м²|шт|дн|ДДУ|₽/м²|[а-яё]{1,4}\s*₽?)?$",
    re.IGNORECASE)


def looks_numeric(text: str) -> bool:
    """Читается ли ячейка числом — для выкладки, а не для графика."""
    raw = str(text or "").strip().replace("−", "-")
    return bool(raw) and bool(_NUMERIC_CELL.match(raw))


# Объявленное отсутствие: на экране это прочерк, и он значит «здесь значения
# нет», а не «здесь ноль».
_BLANK_CELL = {"", "—", "–", "-", "н/д", "нет"}
# Мера колонки — единица, в которой она измерена, и в шапке она стоит где
# придётся: «МЛН ₽, ФАКТ» ставит её впереди, «Факт, млн ₽» — сзади, «Лотов» —
# сама по себе. Брать часть после последней запятой нельзя: на настоящем
# отчёте три денежные колонки получали три разные меры («факт», «план фм»,
# «план банка») и в один график не собирались, а именно они и есть тот самый
# сводный график листа. Единица ищется по своему написанию, а не по месту.
#
# Порядок важен: «₽/м²» проверяется раньше «₽» и раньше «м²», иначе цена метра
# попадёт то в деньги, то в площадь — а это ровно те две оси, между которыми
# её и надо различить.
_MEASURES: tuple[tuple[str, re.Pattern[str]], ...] = (
    ("₽/м²", re.compile(r"(?:₽|руб)\s*/\s*м|цена\s+мет|за\s+метр", re.IGNORECASE)),
    ("млн ₽", re.compile(r"\b(?:млн|млрд|тыс)\b[^,]*(?:₽|руб)|(?:₽|руб)[^,]*\b(?:млн|млрд)\b",
                         re.IGNORECASE)),
    ("₽", re.compile(r"₽|руб", re.IGNORECASE)),
    ("м²", re.compile(r"\bм2\b|м²|метр", re.IGNORECASE)),
    ("%", re.compile(r"%|дол[яию]|конверс", re.IGNORECASE)),
    ("шт", re.compile(r"\b(?:лотов|лоты|штук|шт|договоров|мест|обращений|броней)\b",
                      re.IGNORECASE)),
)


def measure_of(head: str) -> str:
    """В чём измерена колонка. Разные меры на одну ось не кладут."""
    raw = str(head or "").strip()
    for name, marks in _MEASURES:
        if marks.search(raw):
            return name
    return raw.casefold()


def charts(table: dict[str, Any]) -> list[dict[str, Any]]:
    """Графики таблицы: по одному на числовую колонку.

    Класть рубли, метры и цену метра на одну ось нельзя — столбик в пиксель
    рядом со столбиком во весь слайд читается как «этого нет». Первая версия
    решала это выбором одной колонки и подписью, какая нарисована; владелец
    (29.08.2026): «не проще для каждого графика свой слайд сделать?» — проще, и
    ничего не теряется: каждая мера показана, ни одна не спорит с соседней, а
    лишний слайд в PowerPoint удаляют одним нажатием.

    Колонка берётся, только если каждая её ЗАПОЛНЕННАЯ ячейка — число.
    Прочерк — объявленное отсутствие, и он едет пропуском: столбика в этом
    месяце просто нет. Нулём его рисовать нельзя («пропуск в ряду — не ноль»),
    но и выбрасывать всю колонку из-за одного прочерка нельзя тоже — так с
    графика «факт против планов» пропадал план банка, у которого первый
    квартал пустой.
    """
    head, rows = table.get("head") or [], table.get("rows") or []
    if len(head) < 2 or len(rows) < 2:
        return []
    categories = [str(row[0]) for row in rows if row]

    def column(index: int) -> list[float | None] | None:
        values: list[float | None] = []
        for row in rows:
            raw = str(row[index] if index < len(row) else "").strip()
            if raw in _BLANK_CELL:
                values.append(None)
                continue
            number = cell_number(raw)
            if number is None:
                return None
            values.append(float(number))
        if len(values) != len(categories):
            return None
        return values if sum(1 for value in values if value is not None) >= 2 else None

    numeric = {index: column(index) for index in range(1, len(head))}
    numeric = {index: values for index, values in numeric.items() if values}
    # Цена метра — не такая же мера, как метры и рубли: она про другое и живёт
    # линией на своей шкале. «Цена — всегда линия на своей шкале, а не вкладка
    # со столбиками» (владелец, 26.08.2026): на общей шкале с рублями её не
    # видно, а отдельной вкладкой она исчезает ровно тогда, когда смотрят на
    # метры. В колоде она уходила своим слайдом со столбиками — то же самое
    # другими словами. Теперь она идёт линией справа на каждом графике объёма.
    price = next((index for index in numeric if _PRICE.search(str(head[index]))), None)
    line = ({"name": str(head[price]), "values": numeric[price]}
            if price is not None else None)
    # График на раздел — один. Мера на слайд давала три-четыре почти одинаковых
    # столбиковых листа подряд: «столбики не функциональны и не красивы»
    # (владелец, 30.08.2026). Остальные меры не пропадают — они в таблице
    # раздела, которая идёт следом и которую правят.
    #
    # Берётся денежная колонка: управленцу нужны рубли, а штуки и метры при них
    # справочны. Денег нет — первая числовая.
    order = [index for index in numeric if index != price]
    if not order:
        order = list(numeric)
    # Меры, которые экран предлагает переключателем: у динамики это млн ₽, м²
    # и лоты, у планов — млн ₽ и м². В документе переключателя нет, и мера, до
    # которой не переключились, в нём просто отсутствует — то же правило, что у
    # свёрнутой таблицы: раскрыть её читателю нечем. Поэтому лист на меру.
    #
    # Схлопывать их в один график было ошибкой (0.20.83): «куча столбиков»
    # приходила из разделов, где графика нет на экране вовсе, и это лечится
    # признаком `charted`, а не потерей мер.
    groups: list[tuple[str, list[int]]] = []
    for index in order:
        unit = measure_of(str(head[index]))
        found = next((pair for pair in groups if pair[0] == unit), None)
        if found is None:
            groups.append((unit, [index]))
        else:
            found[1].append(index)
    # Денежная мера идёт первой: управленцу нужны рубли, а метры и штуки при
    # них справочны.
    groups.sort(key=lambda pair: 0 if _MONEY.search(pair[0]) else 1)
    money = next((index for index in order if _MONEY.search(str(head[index]))), order[0])
    # Колонки ОДНОЙ меры идут рядами одного графика, а не разъезжаются по
    # слайдам и не теряются. Раздел «Факт против планов» показывал один факт:
    # три колонки в «млн ₽» — факт, план ФМ, план банка, — а на слайд уезжала
    # первая, и график с именем «против планов» никаких планов не показывал.
    # Мера берётся из шапки: «Факт, млн ₽» и «План ФМ, млн ₽» — одна ось,
    # «Лотов», «м²» и «₽/м²» — разные, и класть их вместе нельзя.
    second_columns: list[int] = []
    # Вторая мера идёт линиями на своей шкале справа — так собран сводный
    # график листа: факт столбиками, планы линиями на шкале рублей, все цены
    # линиями на шкале ₽/м². Второй мерой берётся цена, если она есть: «цена —
    # всегда линия на своей шкале» (владелец, 26.08.2026). Третьей меры на
    # графике не бывает — две шкалы, больше некуда.
    price_unit = measure_of(str(head[price])) if price is not None else ""
    if price_unit:
        second_columns = [index for index in numeric
                          if measure_of(str(head[index])) == price_unit]

    drawn: list[dict[str, Any]] = []
    for unit, columns in groups:
        if unit == price_unit:
            continue
        lead = columns[0]
        second = [index for index in second_columns if index not in columns]
        drawn.append({
            "name": str(head[lead]), "categories": categories,
            "values": numeric[lead],
            # Заголовок слайда называет МЕРУ, а не первый ряд: «· Факт, млн ₽»
            # над графиком, где рядом стоят оба плана, обещает меньше, чем на
            # слайде есть. Ряд один — его имя и есть мера.
            "measure": unit if len(columns) > 1 else str(head[lead]),
            # Ряды сверх первого — линиями на той же шкале: столбики в пять
            # рядов читаются частоколом, а план рядом с фактом — линией.
            "extra": [{"name": str(head[index]), "values": numeric[index]}
                      for index in columns[1:]],
            # Цена метра идёт линией на своей шкале на КАЖДОМ графике объёма:
            # «цена — всегда линия на своей шкале» (владелец, 26.08.2026).
            "second": [{"name": str(head[index]), "values": numeric[index]}
                       for index in second],
            "second_measure": price_unit if second else "",
        })
    if not drawn and price is not None:
        # Одна цена без объёма — сама себе график: показать её иначе нечем, а
        # пропустить значило бы потерять единственную меру раздела.
        drawn.append({
            "name": str(head[price]), "categories": categories,
            "values": numeric[price], "measure": str(head[price]),
            "extra": [], "second": [], "second_measure": "",
        })
    return drawn


_SEC_CAT_AX, _SEC_VAL_AX = 771001, 771002


# Цвета рядов. Столбики — фирменный синий; линии первой шкалы и линии второй
# берут свои ряды из палитры листа: два ряда одного цвета неразличимы, а
# радуга читается как разные величины. Порядок повторяет порядок листа.
_PRIMARY_LINES = ("C4581B", "7C6BB5", "8A9BA8")
_SECOND_LINES = ("1F5C87", "D9A441", "4FA07A", "9A6BB5")


def _combo(chart: Any, *, primary: int, secondary: int) -> None:
    """Столбики, линии на своей шкале и линии на второй — как на листе.

    Комбинированных графиков python-pptx не строит, и это единственное место
    модуля, где XML правится руками. Иначе пришлось бы класть цену метра
    столбиками на общую шкалу — рядом с рублями такой столбик выходит в
    пиксель, — либо уносить её отдельным слайдом, что владелец назвал ошибкой:
    цена сравнивается с ценой ТОГО ЖЕ товара, и смотрят на неё вместе с
    объёмом.

    Планы идут линиями на ШКАЛЕ РУБЛЕЙ, а не столбиками рядом: пять столбиков
    в категории читаются частоколом, и «факт против плана» в нём не виден.
    Именно так собран сводный график отчёта, и слайд его повторяет, а не
    придумывает свой.

    Ось цены не от нуля: «урезанная шкала обязана назваться» — она подписана
    справа и её деления видны. Ось объёма остаётся от нуля: у метров и рублей
    ноль — настоящее начало отсчёта.
    """
    if not primary and not secondary:
        return
    from lxml import etree
    from pptx.oxml.ns import qn

    plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    bar = plot_area.find(qn("c:barChart"))
    if bar is None:
        raise DeckUnavailable("столбиков на графике не нашлось")
    series = bar.findall(qn("c:ser"))
    if len(series) < 1 + primary + secondary:
        raise DeckUnavailable("рядов меньше, чем нужно графику листа")
    own_axes = [node.get("val") for node in bar.findall(qn("c:axId"))]

    def to_lines(taken: list[Any], axes: list[str], palette: tuple[str, ...]) -> Any:
        for node in taken:
            bar.remove(node)
            # `invertIfNegative` — свойство столбика; в линии его быть не должно.
            for junk in node.findall(qn("c:invertIfNegative")):
                node.remove(junk)
        # Порядок детей области обязателен: сначала ВСЕ группы графиков, потом
        # оси. Линия, приписанная в конец, встала бы после осей — PowerPoint
        # такой файл не открывает вовсе, а всё остальное его читает и молчит.
        chart_group = etree.Element(qn("c:lineChart"))
        bar.addnext(chart_group)
        etree.SubElement(chart_group, qn("c:grouping")).set("val", "standard")
        etree.SubElement(chart_group, qn("c:varyColors")).set("val", "0")
        for order, node in enumerate(taken):
            chart_group.append(node)
            properties = etree.SubElement(node, qn("c:spPr"))
            stroke = etree.SubElement(properties, qn("a:ln"))
            stroke.set("w", "28575")
            fill = etree.SubElement(stroke, qn("a:solidFill"))
            etree.SubElement(fill, qn("a:srgbClr")).set(
                "val", palette[min(order, len(palette) - 1)])
            etree.SubElement(node, qn("c:smooth")).set("val", "0")
        etree.SubElement(chart_group, qn("c:marker")).set("val", "1")
        for value in axes:
            etree.SubElement(chart_group, qn("c:axId")).set("val", str(value))
        return chart_group

    if secondary:
        to_lines(series[-secondary:], [_SEC_CAT_AX, _SEC_VAL_AX], _SECOND_LINES)
    if primary:
        rest = series[1:len(series) - secondary]
        to_lines(rest, own_axes, _PRIMARY_LINES)

    if not secondary:
        return

    def axis(tag: str, own: int, cross: int, *, position: str, deleted: str):
        node = etree.SubElement(plot_area, qn(tag))
        etree.SubElement(node, qn("c:axId")).set("val", str(own))
        scaling = etree.SubElement(node, qn("c:scaling"))
        etree.SubElement(scaling, qn("c:orientation")).set("val", "minMax")
        etree.SubElement(node, qn("c:delete")).set("val", deleted)
        etree.SubElement(node, qn("c:axPos")).set("val", position)
        etree.SubElement(node, qn("c:crossAx")).set("val", str(cross))
        return node

    price_axis = axis("c:valAx", _SEC_VAL_AX, _SEC_CAT_AX, position="r", deleted="0")
    etree.SubElement(price_axis, qn("c:crosses")).set("val", "max")
    axis("c:catAx", _SEC_CAT_AX, _SEC_VAL_AX, position="b", deleted="1")


def build(pages: list[dict[str, Any]], *, title: str, subtitle: str, footer: str) -> bytes:
    """Колода из разобранных разделов. Ни одного числа здесь не считается."""
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION,
                                     XL_LEGEND_POSITION, XL_TICK_MARK)
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:  # noqa: BLE001
        raise DeckUnavailable(
            "В образе нет python-pptx — презентацию собрать нечем") from exc
    if not pages:
        raise DeckUnavailable("Собирать нечего: в своде нет ни одного раздела")

    deck = Presentation()
    deck.slide_width = Inches(SLIDE_W_IN)
    deck.slide_height = Inches(SLIDE_H_IN)
    blank = deck.slide_layouts[6]
    # Цвета продукта, а не офисные: тот же синий, что в кабинете и на странице.
    # Ряд один, поэтому категориальная палитра здесь не нужна — нужен один
    # фирменный цвет и текстовые токены под подписи.
    # Токены взяты у самого отчёта (`:root` кабинета и его печатные стили), а
    # не подобраны на глаз: у колоды была своя палитра — темнее и глуше, — и
    # рядом с листом она читалась как другой документ. Заголовок раздела в
    # отчёте не чёрный, а приглушённый; столбик светлее фирменного синего.
    ink = RGBColor(0x16, 0x20, 0x2B)      # --ink
    dim = RGBColor(0x5B, 0x6B, 0x7D)      # --dim, он же цвет заголовков разделов
    brand = RGBColor(0x4E, 0x9B, 0xDE)    # столбик графика на листе
    deep = RGBColor(0x13, 0x67, 0xAE)     # --blue
    body = RGBColor(0x33, 0x42, 0x4F)     # текст вывода в плашке
    tile_fill = RGBColor(0xF8, 0xFA, 0xFC)
    note_fill = RGBColor(0xF6, 0xF9, 0xFC)
    paper = RGBColor(0xFF, 0xFF, 0xFF)

    def textbox(slide, text: str, *, top: float, size: int, colour: RGBColor,
                bold: bool = False, height: float = 0.6):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top),
                                       Inches(SLIDE_W_IN - 1.2), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        run = frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        return box

    def spaced(run, hundredths: int) -> None:
        """Разрядка. python-pptx её не знает — ставим атрибут прямо в разметке:
        прописные без разрядки в шапке слипаются в одно слово."""
        run.font._rPr.set("spc", str(hundredths))

    def rule(slide, *, top: float, colour: RGBColor, weight: float = 0.75,
             left: float = 0.6, width: float | None = None) -> None:
        """Волосяная линейка. В PDF ими отбиты и шапка, и колонтитул."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
            Inches(width if width is not None else SLIDE_W_IN - left * 2),
            Pt(weight))
        line.fill.solid()
        line.fill.fore_color.rgb = colour
        line.line.fill.background()
        line.shadow.inherit = False

    def footer_line(slide, number: int, *, name: bool = True) -> None:
        """Колонтитул: лист, отделившийся от колоды, обязан сам говорить, чей
        он и на какую дату. В PDF он повторяется на каждой странице."""
        rule(slide, top=SLIDE_H_IN - 0.62, colour=RGBColor(0xDD, 0xE5, 0xED))
        left = slide.shapes.add_textbox(Inches(0.6), Inches(SLIDE_H_IN - 0.55),
                                        Inches(SLIDE_W_IN - 1.8), Inches(0.35))
        run = left.text_frame.paragraphs[0].add_run()
        run.text = " · ".join(part for part in (footer, title, subtitle) if part) \
            if name else ""
        run.font.size = Pt(9)
        run.font.color.rgb = dim
        right = slide.shapes.add_textbox(Inches(SLIDE_W_IN - 1.4),
                                         Inches(SLIDE_H_IN - 0.55),
                                         Inches(0.8), Inches(0.35))
        paragraph = right.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        page = paragraph.add_run()
        page.text = str(number)
        page.font.size = Pt(9)
        page.font.color.rgb = dim

    def eyebrow_line(slide, text: str, *, top: float) -> float:
        """Надзаголовок печатной шапки: прописные вразрядку под волосяной
        линейкой. В PDF ими открывается документ, и лист колоды открывается
        так же — иначе слайд и страница выглядят двумя разными отчётами."""
        box = textbox(slide, text.upper(), top=top, size=10, colour=dim, height=0.28)
        spaced(box.text_frame.paragraphs[0].runs[0], 85)
        rule(slide, top=top + 0.3, colour=ink)
        return top + 0.42

    def new_slide(heading: str, section: str = ""):
        """Лист раздела в вёрстке печатного отчёта: надзаголовок, линейка,
        заголовок, колонтитул. Номер листа живёт в колонтитуле, а не углом:
        на бумаге он стоит там же."""
        slide = deck.slides.add_slide(blank)
        # Надзаголовок несёт слова отчёта, а не наши: имя раздела, а на его
        # первом листе — подзаголовок самого свода. Придуманная строка на
        # слайде читается как часть отчёта, которой в отчёте нет.
        top = eyebrow_line(slide, section or subtitle or footer, top=0.42)
        textbox(slide, heading, top=top, size=26, colour=dim, bold=True, height=0.75)
        footer_line(slide, len(deck.slides))
        return slide

    def put_table(slide, table: dict[str, Any], *, top: float, height: float) -> None:
        """Таблица в оформлении продукта, а не в заводской синей полосатости.

        Заводской стиль PowerPoint красит шапку в сплошную синеву и чередует
        строки — на слайде с числами это шум, который спорит с числами. Здесь
        шапка на светлой подложке, строки белые, текст носит текстовые токены.
        """
        head = table.get("head") or []
        rows = table.get("rows") or []
        columns = max([len(head)] + [len(row) for row in rows]) or 1
        shape = slide.shapes.add_table(len(rows) + (1 if head else 0), columns,
                                       Inches(0.6), Inches(top),
                                       Inches(SLIDE_W_IN - 1.2), Inches(height))
        grid = shape.table
        grid.first_row = bool(head)
        grid.horz_banding = False

        def dress(cell, text: str, *, header: bool, first: bool) -> None:
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFA) if header else RGBColor(0xFF, 0xFF, 0xFF)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            paragraph = cell.text_frame.paragraphs[0]
            # Числу место справа: колонка чисел, прижатая влево, не читается
            # столбиком. Заголовок стоит там же, где его числа.
            # Первая колонка — имена, и вправо она не уходит никогда: «100%»
            # в ней это условие оплаты, а не число, и прижатое вправо оно
            # встаёт под колонку договоров.
            if not first and (looks_numeric(text) or header):
                paragraph.alignment = PP_ALIGN.RIGHT
            for run in paragraph.runs or [paragraph.add_run()]:
                run.font.size = Pt(11)
                run.font.bold = header
                run.font.color.rgb = ink if header else RGBColor(0x2A, 0x33, 0x3D)

        offset = 0
        if head:
            for index in range(columns):
                dress(grid.cell(0, index),
                      str(head[index]) if index < len(head) else "", header=True,
                      first=index == 0)
            offset = 1
        for line, row in enumerate(rows):
            for index in range(columns):
                dress(grid.cell(line + offset, index),
                      str(row[index]) if index < len(row) else "", header=False,
                      first=index == 0)

    def put_tiles(slide, table: dict[str, Any], *, top: float) -> float:
        """Ключевые числа плитками, как на экране: число крупно, имя над ним.

        Таблицей «Показатель / Значение / Пояснение» они перестают быть тем,
        чем являются: на слайде это первое, на что смотрят, и читаться оно
        должно с трёх метров, а не разбираться построчно.
        """
        rows = [row for row in (table.get("rows") or []) if any(row)]
        if not rows:
            return top
        count = min(len(rows), 4)
        gap, edge = 0.25, 0.6
        width = (SLIDE_W_IN - edge * 2 - gap * (count - 1)) / count
        for index, row in enumerate(rows[:count]):
            left = edge + index * (width + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                          Inches(width), Inches(1.35))
            card.fill.solid()
            card.fill.fore_color.rgb = tile_fill
            card.line.color.rgb = RGBColor(0xDD, 0xE5, 0xED)
            card.shadow.inherit = False
            frame = card.text_frame
            frame.word_wrap = True
            frame.margin_left = frame.margin_right = Inches(0.16)
            frame.margin_top = Inches(0.12)
            for order, (text, size, bold, tone) in enumerate((
                    (str(row[0] if len(row) > 0 else ""), 11, False, dim),
                    (str(row[1] if len(row) > 1 else ""), 24, True, ink),
                    (str(row[2] if len(row) > 2 else ""), 10, False, dim))):
                if not text:
                    continue
                para = frame.paragraphs[0] if order == 0 else frame.add_paragraph()
                para.alignment = PP_ALIGN.LEFT
                run = para.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = tone
        return top + 1.6

    def put_note(slide, text: str, *, top: float, size: int = 14) -> float:
        """Вывод раздела — плашка с синей полосой слева, как в отчёте.

        На листе это `.sumup`: светлая подложка, полоса 3px и текст своим
        цветом. Серой строкой без подложки вывод читается как подпись к
        соседнему, а он и есть ответ раздела.
        """
        if not text:
            return top
        lines = 1 + len(text) // 110
        height = 0.34 + 0.24 * lines
        plate = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(top),
            Inches(SLIDE_W_IN - 1.2), Inches(height))
        plate.fill.solid()
        plate.fill.fore_color.rgb = note_fill
        plate.line.fill.background()
        plate.shadow.inherit = False
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(top), Pt(3), Inches(height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = brand
        bar.line.fill.background()
        bar.shadow.inherit = False
        box = slide.shapes.add_textbox(Inches(0.78), Inches(top + 0.06),
                                       Inches(SLIDE_W_IN - 1.56), Inches(height - 0.12))
        frame = box.text_frame
        frame.word_wrap = True
        run = frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = body
        return top + height + 0.18

    def put_shelf(slide, table: dict[str, Any], *, top: float) -> float:
        """Полка показателей титула — та же, что в шапке печатного отчёта:
        одна подложка, равные колонки, разделённые волосяными линейками.

        Пять чисел вразброс по белому читаются как обрывки текста, а на своей
        плашке — как одна панель. Колонки равные и разделены линейками, иначе
        длинная сноска второго столбца перекашивает весь ряд.
        """
        rows = [row for row in (table.get("rows") or []) if any(row)][:5]
        if not rows:
            return top
        left, width, height = 0.6, SLIDE_W_IN - 1.2, 1.25
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                       Inches(top), Inches(width), Inches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = tile_fill
        panel.line.color.rgb = RGBColor(0xE3, 0xEB, 0xF2)
        panel.shadow.inherit = False
        panel.adjustments[0] = 0.04
        column = width / len(rows)
        for index, row in enumerate(rows):
            if index:
                divider = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(left + index * column),
                    Inches(top + 0.16), Pt(0.75), Inches(height - 0.32))
                divider.fill.solid()
                divider.fill.fore_color.rgb = RGBColor(0xE3, 0xEB, 0xF2)
                divider.line.fill.background()
                divider.shadow.inherit = False
            box = slide.shapes.add_textbox(Inches(left + index * column + 0.14),
                                           Inches(top + 0.14),
                                           Inches(column - 0.28), Inches(height - 0.28))
            frame = box.text_frame
            frame.word_wrap = True
            # Порядок печатной плитки: число, под ним имя, под ним сноска.
            # На экране имя стоит над числом, на бумаге — под: там первым
            # смотрят на само число.
            # Порядок плитки — как в отчёте: имя, под ним число, под ним
            # сноска. Число сверху я поставил сам, и это была не та плитка.
            for order, (text, size, bold, tone) in enumerate((
                    (str(row[0] if len(row) > 0 else ""), 11, False, dim),
                    (str(row[1] if len(row) > 1 else ""), 20, True, ink),
                    (str(row[2] if len(row) > 2 else ""), 9, False,
                     RGBColor(0x7B, 0x8B, 0x9A)))):
                if not text:
                    continue
                para = frame.paragraphs[0] if order == 0 else frame.add_paragraph()
                para.alignment = PP_ALIGN.LEFT
                run = para.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = tone
        return top + height + 0.2

    def put_strip(slide, strip: dict[str, Any], *, top: float) -> float:
        """Лента долей — теми же цветами, что на экране, и с подписями.

        На слайде от неё не оставалось ничего: у кусков ленты нет текста, одна
        ширина и цвет. А это и есть ответ раздела — как устроен пул и как из
        него покупают.
        """
        parts = [item for item in strip.get("parts") or [] if item.get("share")]
        if not parts:
            return top
        caption = str(strip.get("caption") or "")
        if caption:
            textbox(slide, caption, top=top, size=12, colour=dim, height=0.3)
            top += 0.32
        total = sum(item["share"] for item in parts) or 100.0
        width = SLIDE_W_IN - 1.2
        left = 0.6
        narrow = []
        for item in parts:
            span = width * item["share"] / total
            block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                           Inches(max(span, 0.02)), Inches(0.62))
            block.fill.solid()
            block.fill.fore_color.rgb = RGBColor.from_string(item["colour"])
            block.line.fill.background()
            block.shadow.inherit = False
            name = str(item["name"])
            if span >= 1.3:
                # Имя стоит в своём куске: серой строкой под лентой оно с
                # цветом не связано вовсе, а подсказки по наведению на слайде
                # не бывает.
                para = block.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = name
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = paper
            else:
                narrow.append(name)
            left += span
        top += 0.70
        if narrow:
            # Узкому куску имя внутрь не влезает — оно уходит строкой под
            # ленту, но только оно: остальные уже подписаны.
            textbox(slide, "Узкие полосы: " + "; ".join(narrow),
                    top=top, size=10, colour=dim, height=0.3)
            top += 0.34
        return top

    def put_chart(slide, data: dict[str, Any], *, top: float, height: float) -> None:
        """График, а не заводская заготовка PowerPoint.

        «Убогие графики очень, как для первого класса школы» (владелец,
        30.08.2026) — и это была правда: столбики во всю ширину слота заводской
        синевы, сетка по всему полю, ось со значениями и подписи офисным
        шрифтом. Здесь наведён порядок по правилам оформления данных:
        столбик тонкий, цвет — наш фирменный, а не офисный; текст носит
        текстовые токены, а не цвет ряда; сетка убрана там, где значения стоят
        прямо на столбиках, и оставлена волосяной, где их много; легенды нет —
        ряд один, и его называет заголовок слайда.
        """
        extra = list(data.get("extra") or [])
        second = list(data.get("second") or [])
        payload = CategoryChartData()
        payload.categories = data["categories"]
        payload.add_series(data["name"], data["values"])
        for other in extra + second:
            payload.add_series(other["name"], other["values"])
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(top),
            Inches(SLIDE_W_IN - 1.2), Inches(height), payload)
        chart = frame.chart
        # Своего заголовка у графика нет: его имя уже стоит заголовком слайда,
        # а повторённое мелким серым оно читается как чужая подпись.
        chart.has_title = False
        # Ряд один — легенды нет, его называет заголовок слайда. Рядов
        # несколько — без легенды они неразличимы.
        chart.has_legend = bool(extra or second)
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.font.size = Pt(11)
        chart.font.name = "Calibri"
        chart.font.color.rgb = dim

        plot = chart.plots[0]
        plot.vary_by_categories = False
        # Столбик не заполняет слот: воздух между столбиками — часть чтения, а
        # столбик во всю ширину слота и есть то самое «как для первого класса».
        # Ширины в формате нет, есть просвет — и он тем больше, чем меньше
        # категорий: на трёх категориях столбик иначе выходит в две ладони.
        count = len(data["categories"])
        plot.gap_width = 400 if count <= 3 else (250 if count <= 8 else
                                                 (140 if count <= 12 else 60))

        # Ряды одной меры — оттенками одного цвета, а не радугой: они про одно
        # и то же, и разный цвет читался бы как разные величины. Факт носит
        # фирменный, планы — бледнее: смотрят на факт.
        tones = (brand, RGBColor(0x7F, 0xB2, 0xE5), RGBColor(0xB9, 0xCF, 0xE4),
                 RGBColor(0xD7, 0xE4, 0xF0))
        # Столбики только у первого ряда: остальные уедут линиями. Пять
        # столбиков в категории читаются частоколом, и «факт против плана» в
        # нём не виден — на листе планы идут линиями, и слайд повторяет лист.
        bars = plot.series[0]
        bars.format.fill.solid()
        bars.format.fill.fore_color.rgb = tones[0]
        bars.format.line.fill.background()

        # Число на каждом столбике — это хаос, если столбиков много: тогда
        # значения несёт ось. Мало — значения стоят прямо на шапках, и ось со
        # своей сеткой становится лишней краской.
        # Со второй шкалой первая обязана остаться видимой. Удалённая, она
        # уводит столбики на шкалу цены — 35 против 800 000, — и на слайде от
        # них не остаётся ничего: «тут просто линии» (владелец, 30.08.2026).
        has_line = bool(second)
        # Значения стоят на столбиках, пока их не много. Ряд столбиков один,
        # поэтому считаются его категории — линии подписей не несут.
        labelled = len(data["categories"]) <= 8
        plot.has_data_labels = labelled
        if labelled:
            labels = plot.data_labels
            labels.number_format = "#,##0.#"
            labels.number_format_is_linked = False
            labels.position = XL_LABEL_POSITION.OUTSIDE_END
            labels.font.size = Pt(11)
            labels.font.bold = True
            labels.font.color.rgb = ink

        value_axis = chart.value_axis
        value_axis.has_major_gridlines = not labelled and not has_line
        if value_axis.has_major_gridlines:
            # Имя `line` здесь занято рядом цены — сетка носит своё.
            hairline = value_axis.major_gridlines.format.line
            hairline.color.rgb = RGBColor(0xE3, 0xEB, 0xF2)
            hairline.width = Pt(0.75)
        value_axis.visible = has_line or not labelled
        value_axis.has_minor_gridlines = False
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.format.line.fill.background()
        if value_axis.visible:
            value_axis.tick_labels.number_format = "#,##0"
            value_axis.tick_labels.number_format_is_linked = False
            value_axis.tick_labels.font.size = Pt(10)
            value_axis.tick_labels.font.color.rgb = dim

        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.format.line.color.rgb = RGBColor(0xDD, 0xE5, 0xED)
        category_axis.tick_labels.font.size = Pt(10)
        category_axis.tick_labels.font.color.rgb = dim
        _combo(chart, primary=len(extra), secondary=len(second))

    # Титул: чей отчёт и на какую дату. Слайд, отделившийся от колоды, обязан
    # сам говорить, чей он, — как лист на бумаге.
    # Титул — печатная шапка отчёта: надзаголовок под линейкой, крупное имя,
    # подзаголовок. «Максимально похожий на PDF» (владелец, 31.08.2026), а PDF
    # у нас белый — тёмный лист ему противоречил.
    first = deck.slides.add_slide(blank)
    top = eyebrow_line(first, footer, top=1.35)
    textbox(first, title, top=top, size=32, colour=ink, bold=True, height=0.95)
    top += 1.05
    if subtitle:
        # Подзаголовок печатного отчёта: одна строка о том, на чём посчитано.
        textbox(first, subtitle, top=top, size=14, colour=dim, height=0.4)
        top += 0.5
    # Полка показателей стоит ДО первого графика — как в отчёте: главные числа
    # не вылавливают из картинки. Плитки берутся у того раздела, где они есть,
    # а не считаются заново: второй счёт той же величины однажды разойдётся.
    shelf = next((table for page in pages for table in (page.get("tables") or [])
                  if table.get("kind") == "tiles" and any(
                      any(row) for row in (table.get("rows") or []))), None)
    if shelf is not None:
        # Помечаем взятую полку: повторённая через слайд, она читается как
        # вторые числа о том же.
        shelf["used"] = True
        put_shelf(first, shelf, top=top + 0.1)

    # На титуле имя отчёта уже стоит заголовком: повторённое колонтитулом,
    # оно читается как заводская рамка. Номер листа остаётся.
    footer_line(first, 1, name=False)

    # Раздел, у которого показать нечего, кроме вывода, листом не является:
    # «слайды 2-3 пустые вообще, там по одной строчке текста» (владелец,
    # 31.08.2026). На листе такой раздел занимает три сантиметра в потоке, а
    # на слайде — заголовок и пять дюймов белого. Такие идут по нескольку на
    # общий лист, в своём порядке и со своими заголовками.
    THIN_PER_SLIDE = 3

    def thin(page: dict[str, Any]) -> bool:
        return not (page.get("tables") or page.get("strips") or page.get("lines"))

    pending: list[dict[str, Any]] = []

    def flush_thin() -> None:
        while pending:
            batch, pending[:] = pending[:THIN_PER_SLIDE], pending[THIN_PER_SLIDE:]
            slide = new_slide(str(batch[0].get("title") or title))
            top = CONTENT_TOP
            for index, item in enumerate(batch):
                if index:
                    textbox(slide, str(item.get("title") or ""), top=top, size=17,
                            colour=dim, bold=True, height=0.4)
                    top += 0.5
                top = put_note(slide, str(item.get("note") or ""), top=top)
                top += 0.15

    for page in pages:
        # У шапки свода своего заголовка нет — она несёт ключевые числа
        # проекта. «Раздел» над плитками не говорит ничего.
        heading = str(page.get("title") or title)
        tables = list(page.get("tables") or [])
        note = str(page.get("note") or "").strip()
        lines = [line for line in (page.get("lines") or []) if line][:LINES_PER_SLIDE]

        # Первый слайд раздела — о чём он: вывод крупно и подписи, которые на
        # экране стоят рядом с картинкой (доли полос — это и есть числа такого
        # раздела). Слайда НЕТ, когда класть на него нечего: «Расторжения» с
        # одним заголовком и пустым полем — это не раздел, а пустой лист.
        strips = list(page.get("strips") or [])
        tiles = [table for table in tables
                 if table.get("kind") == "tiles" and not table.get("used")]
        tables = [table for table in tables if table.get("kind") != "tiles"]
        if not (tables or strips or tiles or lines) and note:
            # Копим и кладём по нескольку на лист — своего листа такой раздел
            # не заслуживает, но и пропасть не должен: вывод и есть его ответ.
            pending.append(page)
            continue
        flush_thin()
        # График на слайде — только там, где он есть на экране. Прежде колода
        # заводила столбики под ПЕРВУЮ таблицу каждого раздела, и на своде из
        # десяти разделов выходило восемь почти одинаковых синих слайдов —
        # «там была куча столбиков опять» (владелец, 31.08.2026). Расторжения
        # столбиками не читаются вовсе, у структуры оплаты их три штуки, а
        # экран в этих разделах и не рисует ничего: он рисует график ровно в
        # трёх — динамика, факт против планов, эскроу. Это решение уже
        # принято, и второй раз его принимать нельзя: колода взяла бы на себя
        # выбор, которого экран ей не поручал.
        #
        # Лента долей столбиков не получает по той же причине: она говорит то
        # же самое и теми же цветами, что на экране. Числа при этом не
        # пропадают — они в таблице раздела.
        drawn = charts(tables[0]) if (page.get("charted") and tables and not strips) else []
        # Слайд заводится, только если на нём есть что показать. Вывод сам по
        # себе слайдом не является: «этот слайд странный» (владелец,
        # 30.08.2026) — заголовок, одна строка и пять дюймов белого. Такой
        # вывод едет подзаголовком на первый слайд раздела, где есть картинка.
        rich = bool(lines or tiles or strips or (not drawn and tables))
        carry = "" if rich else note

        def lead(slide, top: float) -> float:
            """Вывод раздела над содержимым — один раз на раздел."""
            nonlocal carry
            if not carry:
                return top
            top = put_note(slide, carry, top=top)
            carry = ""
            return top

        opening = None
        if rich:
            opening = new_slide(heading)
            top = CONTENT_TOP
            if tiles:
                top = put_tiles(opening, tiles[0], top=top)
            for strip in strips:
                top = put_strip(opening, strip, top=top + 0.1)
            if lines and not strips:
                # Больше пяти подписей — в две колонки: столбик в двадцать
                # строк уезжает за нижний край, а половина листа стоит пустой.
                columns = 2 if len(lines) > 5 else 1
                per = -(-len(lines) // columns)
                width = (SLIDE_W_IN - 1.2) / columns - 0.2
                for index, line in enumerate(lines):
                    box = opening.shapes.add_textbox(
                        Inches(0.6 + (index // per) * (width + 0.2)),
                        Inches(top + (index % per) * 0.42),
                        Inches(width), Inches(0.4))
                    box.text_frame.word_wrap = True
                    run = box.text_frame.paragraphs[0].add_run()
                    run.text = line
                    run.font.size = Pt(14)
                    run.font.color.rgb = ink
            # Таблица раздела встаёт на тот же лист, если помещается. На
            # экране лента и её числа стоят одним блоком, и лист с одной
            # короткой лентой над пятью дюймами белого — это не «просторно», а
            # «здесь ничего нет». Не поместилась — уезжает своим листом
            # целиком, а не ужимается до нечитаемого.
            if not drawn and tables and not tiles:
                height = 0.34 * (len(tables[0]["rows"]) + 1)
                if top + height <= SLIDE_H_IN - 1.5:
                    put_table(opening, tables[0], top=top + 0.15, height=height)
                    top += 0.15 + height
                    tables = tables[1:]
            # Вывод раздела стоит ПОД его содержимым, как на листе: наверху он
            # читается как подпись к заголовку, а он — ответ раздела.
            if note:
                put_note(opening, note, top=min(top + 0.2, SLIDE_H_IN - 1.35), size=15)
                note = ""

        # По слайду на график: каждая мера показана и ни одна не спорит с
        # соседней. Ряд один, поэтому легенда не нужна — мера стоит в
        # заголовке слайда. Сноски под каждым графиком больше нет: повторённая
        # двадцать раз, она перестаёт быть пояснением и становится шумом —
        # сказать это достаточно один раз, на титуле.
        for chart in drawn:
            part = new_slide(
                f"{heading} · {chart.get('measure') or chart['name']}", heading)
            top = lead(part, CONTENT_TOP)
            put_chart(part, chart, top=top, height=SLIDE_H_IN - top - 0.8)

        # Таблицы — целиком и ячейками: их и правят. Длинная продолжается
        # следующим слайдом, а не ужимается до нечитаемого.
        for table in tables:
            rows = table.get("rows") or []
            for start in range(0, len(rows), ROWS_PER_SLIDE):
                chunk = {"head": table.get("head") or [],
                         "rows": rows[start:start + ROWS_PER_SLIDE]}
                part = new_slide(heading + ("" if start == 0 else " · продолжение"),
                                 heading)
                top = lead(part, CONTENT_TOP)
                put_table(part, chunk, top=top,
                          height=min(SLIDE_H_IN - top - 0.8,
                                     0.34 * (len(chunk["rows"]) + 1)))
        if carry:
            # Разделу нечего показать, кроме вывода: тогда он и есть слайд.
            put_note(new_slide(heading), carry, top=CONTENT_TOP, size=18)
        if note and opening is not None:
            opening.notes_slide.notes_text_frame.text = note

    flush_thin()

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def file_name(title: str) -> str:
    """Имя файла: то же правило, что у PDF отчёта."""
    keep = "".join(ch for ch in title if ch.isalnum() or ch in " -_")[:80].strip()
    return re.sub(r"\s+", " ", keep) or "sales"
